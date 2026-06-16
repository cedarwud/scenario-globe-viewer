#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Build the 繁體中文 requirement-coverage academic presentation (.pptx).

Phase 2 of the requirement-presentation work. Reads its content from
(authoritative, already committed):
  - docs/requirement-presentation-walkthrough.md   (34-row per-requirement)
  - deliverable/slide-content-inventory.md          (catalog [1]-[16], claim bank)
  - docs/data-source-index.md                       (honesty contract)

Honesty contract enforced here: model / standards-derived / geometric-derived /
assumed-Tier-B / public-derived values are NEVER labeled measured or
operator-validated. External gaps (packet test, native RF handover, operator RF
hardware, external acceptance threshold) are shown only as "已揭露待補 / 資料缺口".

Construction: python-pptx on a copy of template.pptx (16:9 academic theme).
The 3 sample slides are removed; masters/layouts/theme are kept. Charts are
generated with matplotlib (CJK = Noto Sans TC). Screenshots are reused from
retained output/ artifacts. Re-runnable: `python3 scripts/build_requirement_presentation.py`.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import demo_scenario as ds   # demo window/route single source of truth

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(ROOT, "template.pptx")
CHART_DIR = os.path.join(ROOT, "output", "presentation-charts")
OUT = os.path.join(ROOT, "deliverable", "requirement-presentation-2026-06-14.pptx")

# ---------------------------------------------------------------- palette
INK      = RGBColor(0x1A, 0x22, 0x38)   # deep navy text
DARKBG   = RGBColor(0x12, 0x19, 0x2B)   # title / section / closing bg
PAPER    = RGBColor(0xFF, 0xFF, 0xFF)   # content bg
SOFT     = RGBColor(0x5B, 0x63, 0x72)   # muted gray
HAIR     = RGBColor(0xD9, 0xDD, 0xE3)   # hairline
BRAND    = RGBColor(0x2A, 0x9D, 0x8F)   # teal accent / brand
BRANDLT  = RGBColor(0x8F, 0xD4, 0xCC)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# status
ST = {
    "complete": (RGBColor(0x2E, 0x7D, 0x32), "完成 complete"),
    "partial":  (RGBColor(0xE0, 0x8A, 0x1E), "部分 partial · Tier B"),
    "gap":      (RGBColor(0xC6, 0x28, 0x28), "資料缺口 source gap"),
    "scope":    (RGBColor(0x6B, 0x72, 0x80), "範圍外 out-of-scope"),
}
# source classes -> (繁中, english tag, color)
CLS = {
    "public":    ("公開來源", "Public source",      RGBColor(0x15, 0x65, 0xC0)),
    "standards": ("標準模型推導", "Standards-derived", RGBColor(0x1B, 0x80, 0x74)),
    "geometric": ("幾何推導", "Geometric-derived",  RGBColor(0x6A, 0x4C, 0x93)),
    "assumed":   ("假設參數", "Assumed-Tier-B",     RGBColor(0xC9, 0x6A, 0x00)),
    "display":   ("顯示轉換", "Display transform",  RGBColor(0x55, 0x6B, 0x78)),
    "external":  ("外部證據", "External evidence",  RGBColor(0x37, 0x47, 0x4F)),
    "legacy":    ("舊快取", "Legacy cache",         RGBColor(0x8D, 0x6E, 0x63)),
    "gap":       ("資料缺口", "Source gap",         RGBColor(0xC6, 0x28, 0x28)),
}
BUCKET_COLOR = {"A": BRAND, "B": RGBColor(0xE0, 0x8A, 0x1E), "C": RGBColor(0x6B, 0x72, 0x80)}

EA_FONT = "Microsoft JhengHei"   # viewer-side 繁中 default
LAT_FONT = "Arial"

SW, SH = Inches(13.333), Inches(7.5)
FOOT = "需求對應簡報 · 2026-06-14 · 來源誠實合約：docs/data-source-index.md"

# screenshots (retained, neutral selected-pair)
def shot(name):
    p = os.path.join(ROOT, "output", name)
    return p if os.path.exists(p) else None

IMG_GLOBE   = shot("playwright/scenario-globe-viewer-selected-pair-1280x800.png")
IMG_LINK    = shot("playwright/selected-pair-3lane-link-map-1280x800.png")
IMG_HANDOVER= shot("playwright/selected-pair-3lane-handover-1280x800.png")
IMG_VIS     = shot("playwright/selected-pair-visible-link-map-1280x800.png")
IMG_AUDIT   = shot("playwright/selected-pair-audit-1280x800.png")
IMG_REPORT  = shot("selected-pair-source-report/source-report-panel-open.png")
IMG_DRAWER  = shot("selected-pair-source-report/source-report-drawer-open.png")

def ch(name):
    return os.path.join(CHART_DIR, name)

# ============================================================ charts
def build_charts():
    os.makedirs(CHART_DIR, exist_ok=True)
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib import font_manager as fm
    for fp in ("/home/u24/.fonts/NotoSansTC-VF.ttf",):
        if os.path.exists(fp):
            try: fm.fontManager.addfont(fp)
            except Exception: pass
    plt.rcParams["font.family"] = ["Noto Sans TC", "WenQuanYi Zen Hei", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False
    teal="#2A9D8F"; amber="#E08A1E"; gray="#6B7280"; navy="#1A2238"; purple="#6A4C93"; blue="#1565C0"

    # 1) bucket distribution
    fig, ax = plt.subplots(figsize=(6.2, 3.0), dpi=150)
    labels=["A 專案 demo (19)", "B Tier B 替代 (7)", "C 外部/報告層 (8)"]
    vals=[19,7,8]; cols=[teal,amber,gray]
    b=ax.barh(labels, vals, color=cols)
    ax.invert_yaxis(); ax.set_xlim(0,21)
    for i,v in enumerate(vals): ax.text(v+0.3, i, str(v), va="center", fontsize=12, color=navy, fontweight="bold")
    ax.set_title("需求分布（共 34 列）", fontsize=13, color=navy, fontweight="bold")
    ax.spines[["top","right"]].set_visible(False); ax.tick_params(labelsize=10)
    fig.tight_layout(); fig.savefig(ch("chart_buckets.png"), bbox_inches="tight", pad_inches=0.14); plt.close(fig)

    # 2) loaded actors
    fig, ax = plt.subplots(figsize=(6.2, 3.0), dpi=150)
    cats=["OneWeb LEO","Galileo MEO","GEO 群組"]; vals=[651,33,574]
    bars=ax.bar(cats, vals, color=[teal,purple,blue])
    for r,v in zip(bars,vals): ax.text(r.get_x()+r.get_width()/2, v+8, str(v), ha="center", fontsize=12, color=navy, fontweight="bold")
    ax.set_ylim(0,730); ax.set_ylabel("actor 數", fontsize=10)
    ax.set_title("來源衛星 actor（共 1258；2026-06-13 快照；runtime 截頂 200/100/60）", fontsize=11, color=navy, fontweight="bold")
    ax.spines[["top","right"]].set_visible(False); ax.tick_params(labelsize=10)
    ax.text(0.99,0.95,"≥500 LEO 由 OneWeb 滿足", transform=ax.transAxes, ha="right", va="top", fontsize=9, color="#555")
    fig.tight_layout(); fig.savefig(ch("chart_actors.png"), bbox_inches="tight", pad_inches=0.14); plt.close(fig)

    # 3) antenna combined gain
    fig, ax = plt.subplots(figsize=(6.2, 3.2), dpi=150)
    cats=["LEO","MEO","GEO"]; vals=[75.7,83.9,91.9]
    bars=ax.bar(cats, vals, color=[teal,amber,blue])
    for r,v in zip(bars,vals): ax.text(r.get_x()+r.get_width()/2, v+0.6, f"{v} dB", ha="center", fontsize=11, color=navy, fontweight="bold")
    ax.set_ylim(0,100); ax.set_ylabel("合併增益 (dB)", fontsize=10)
    ax.set_title("合併天線增益（S.1528 / S.465 模型 + 假設參數）", fontsize=12, color=navy, fontweight="bold")
    ax.spines[["top","right"]].set_visible(False); ax.tick_params(labelsize=10)
    ax.text(0.5,-0.30,"離軸損 2.1–2.3 dB · RSRP 為相對 proxy（EIRP 未知）· 標準模型推導 + 假設參數 · 非營運商硬體",
            transform=ax.transAxes, ha="center", va="top", fontsize=7.8, color="#C0392B")
    fig.tight_layout(); fig.savefig(ch("chart_antenna.png"), bbox_inches="tight", pad_inches=0.14); plt.close(fig)

    # 4) rain attenuation anchors
    fig, ax = plt.subplots(figsize=(6.2, 3.2), dpi=150)
    cats=["緯度 0°","demo 45°","緯度 78°"]; vals=[30.5,27.3,8.3]
    bars=ax.bar(cats, vals, color=["#1565C0","#2A9D8F","#8FD4CC"])
    for r,v in zip(bars,vals): ax.text(r.get_x()+r.get_width()/2, v+0.5, f"{v} dB", ha="center", fontsize=11, color=navy, fontweight="bold")
    ax.set_ylim(0,34); ax.set_ylabel("雨衰 A (dB)", fontsize=10)
    ax.set_title("降雨衰減模型錨點（P.618-14；20 GHz / 50 mm/h）", fontsize=12, color=navy, fontweight="bold")
    ax.spines[["top","right"]].set_visible(False); ax.tick_params(labelsize=10)
    ax.text(0.5,-0.30,"標準模型推導 · 範例 what-if 輸入 · 非實測天氣；晴空報告 rain=0",
            transform=ax.transAxes, ha="center", va="top", fontsize=8.0, color="#C0392B")
    fig.tight_layout(); fig.savefig(ch("chart_rain.png"), bbox_inches="tight", pad_inches=0.14); plt.close(fig)
    print("charts ->", CHART_DIR)

# ============================================================ pptx helpers
def _set_font(run, latin=LAT_FONT, ea=EA_FONT):
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "zh-TW")
    rPr.set("altLang", "zh-TW")
    rPr.set("noProof", "1")
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", ea if tag == "a:ea" else latin)

def bg(slide, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()
    r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element); slide.shapes._spTree.insert(2, r._element)
    return r

def rect(slide, l, t, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    r = slide.shapes.add_shape(shape, l, t, w, h)
    if color is None:
        r.fill.background()
    else:
        r.fill.solid(); r.fill.fore_color.rgb = color
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line; r.line.width = Pt(0.75)
    r.shadow.inherit = False
    return r

def tb(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    b = slide.shapes.add_textbox(l, t, w, h); tf = b.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    return tf

def para(tf, runs, size=13, align=PP_ALIGN.LEFT, space_after=4, space_before=0,
         line=None, first=False, bullet=False, level=0):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.level = level
    if space_after is not None: p.space_after = Pt(space_after)
    if space_before is not None: p.space_before = Pt(space_before)
    if line is not None:
        p.line_spacing = line
    if not bullet:
        pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
        bn = pPr.makeelement(qn("a:buNone"), {});
        for ex in pPr.findall(qn("a:buNone")): pPr.remove(ex)
        pPr.append(bn)
    for item in runs:
        txt, opts = item if isinstance(item, tuple) else (item, {})
        r = p.add_run(); r.text = txt
        r.font.size = Pt(opts.get("sz", size))
        r.font.bold = opts.get("b", False)
        r.font.italic = opts.get("i", False)
        r.font.color.rgb = opts.get("c", INK)
        _set_font(r, ea=opts.get("ea", EA_FONT))
    return p

def footer(slide, idx, total, dark=False):
    c = RGBColor(0x9A, 0xA3, 0xAF) if dark else SOFT
    t = tb(slide, Inches(0.45), Inches(7.04), Inches(10.6), Inches(0.36))
    para(t, [(FOOT, {"sz": 8.5, "c": c})], size=8.5, first=True, space_after=0)
    n = tb(slide, Inches(11.4), Inches(7.04), Inches(1.5), Inches(0.36))
    para(n, [(f"{idx} / {total}", {"sz": 9, "c": c})], align=PP_ALIGN.RIGHT, first=True, space_after=0)

def pic_framed(slide, path, l, t, w, caption=None):
    if not path or not os.path.exists(path): return None
    p = slide.shapes.add_picture(path, l, t, width=w)
    fr = rect(slide, l, t, w, p.height, None, line=HAIR)
    if caption:
        ct = tb(slide, l, t + p.height + Inches(0.04), w, Inches(0.34))
        para(ct, [(caption, {"sz": 8.6, "c": SOFT, "i": True})], first=True, space_after=0)
    return p

# ============================================================ slide builders
def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def slide_title(prs):
    s = add_blank(prs); bg(s, DARKBG)
    rect(s, 0, Inches(2.55), SW, Inches(0.06), BRAND)
    t = tb(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(1.4))
    para(t, [("LEO / MEO / GEO 衛星地面站場景檢視器", {"sz": 25, "b": True, "c": BRANDLT})], first=True, space_after=4)
    t2 = tb(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(1.5))
    para(t2, [("需求逐條對應 ｜ 資料來源誠實合約", {"sz": 33, "b": True, "c": WHITE})], first=True, space_after=6)
    para(t2, [("34 項需求 × 「如何滿足 · 資料來源 · 引用 · 邊界 · 狀態」", {"sz": 15, "c": RGBColor(0xC7,0xCE,0xDA)})], space_after=0)
    leg = tb(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(2.0))
    para(leg, [("資料來源分級：", {"sz": 13, "b": True, "c": WHITE})], first=True, space_after=6)
    para(leg, [("公開來源 · 標準模型推導 · 幾何推導 · 假設參數(Tier-B) · 顯示轉換 · 外部證據 · 舊快取 · 資料缺口",
                {"sz": 12.5, "c": RGBColor(0xC7,0xCE,0xDA)})], space_after=8)
    para(leg, [("誠實鐵則：", {"sz": 13, "b": True, "c": RGBColor(0xF1,0xB0,0x70)}),
               ("模型 / 標準 / 幾何 / 假設 / 公開推導之值，絕不標示為「實測」或「營運商驗證」。",
                {"sz": 12.5, "c": RGBColor(0xE7,0xC9,0xA8)})], space_after=4)
    para(leg, [("外部缺口（封包測試、原生 RF 換手、營運商 RF 硬體、外部驗收門檻）僅呈現為「已揭露待補」。",
                {"sz": 12.5, "c": RGBColor(0xE7,0xC9,0xA8)})], space_after=0)
    d = tb(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.5))
    para(d, [("狀態日 2026-06-14 ｜ 對齊 merge 1496a03（2026-06-13 公開資料缺口關閉）", {"sz": 11, "c": RGBColor(0x8A,0x93,0xA6)})], first=True, space_after=0)
    return s

def slide_section(prs, idx, total, kicker, title, subtitle, color):
    s = add_blank(prs); bg(s, DARKBG)
    rect(s, Inches(0.9), Inches(2.7), Inches(0.9), Inches(0.10), color)
    t = tb(s, Inches(0.9), Inches(2.05), Inches(11.5), Inches(0.6))
    para(t, [(kicker, {"sz": 15, "b": True, "c": color})], first=True, space_after=0)
    t2 = tb(s, Inches(0.9), Inches(2.95), Inches(11.5), Inches(1.2))
    para(t2, [(title, {"sz": 34, "b": True, "c": WHITE})], first=True, space_after=4)
    if subtitle:
        st = tb(s, Inches(0.9), Inches(4.25), Inches(11.3), Inches(2.0))
        for ln in subtitle:
            para(st, [(ln, {"sz": 14, "c": RGBColor(0xC7,0xCE,0xDA)})], first=(ln is subtitle[0]), space_after=6, line=1.1)
    footer(s, idx, total, dark=True)
    return s

def header(slide, title_runs, subtitle=None):
    rect(slide, 0, 0, Inches(0.16), SH, slide._accent)
    rect(slide, Inches(0.45), Inches(1.18), Inches(12.45), Pt(1.0), HAIR)
    t = tb(slide, Inches(0.45), Inches(0.34), Inches(9.7), Inches(0.85), anchor=MSO_ANCHOR.MIDDLE)
    para(t, title_runs, first=True, space_after=0, line=1.02)
    if subtitle:
        st = tb(slide, Inches(0.47), Inches(0.92), Inches(9.7), Inches(0.3))
        para(st, [(subtitle, {"sz": 11, "c": SOFT})], first=True, space_after=0)

def status_pill(slide, status):
    color, label = ST[status]
    r = rect(slide, Inches(10.35), Inches(0.40), Inches(2.55), Inches(0.5), color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = r.text_frame; tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    para(tf, [(label, {"sz": 11.5, "b": True, "c": WHITE})], align=PP_ALIGN.CENTER, first=True, space_after=0)

def class_runs(keys):
    runs = []
    for k in keys:
        zh, en, col = CLS[k]
        runs.append((f"〔{zh} ", {"sz": 12, "b": True, "c": col}))
        runs.append((f"{en}〕", {"sz": 10.5, "b": True, "c": col}))
    return runs

def slide_req(prs, idx, total, r):
    s = add_blank(prs); bg(s, PAPER)
    s._accent = BUCKET_COLOR[r["bucket"]]
    title_runs = [(r["id"] + "　", {"sz": 21, "b": True, "c": BRAND if r["bucket"]=="A" else s._accent}),
                  (r["name"], {"sz": 20, "b": True, "c": INK})]
    header(s, title_runs)
    status_pill(s, r["status"])
    bucket_tag = {"A": "Bucket A · 專案 demo（幾何 / 標準推導，顯示轉換）",
                  "B": "Bucket B · Tier B 公開標準替代（待 Tier A 營運商資料）",
                  "C": "Bucket C · 外部基礎設施 / 報告層（不在 demo 範圍）"}[r["bucket"]]
    bt = tb(s, Inches(0.47), Inches(0.95), Inches(9.6), Inches(0.26))
    para(bt, [(bucket_tag, {"sz": 10.5, "b": True, "c": s._accent})], first=True, space_after=0)

    has_img = bool(r.get("img") or r.get("chart"))
    body_w = Inches(7.35) if has_img else Inches(12.1)
    body = tb(s, Inches(0.5), Inches(1.42), body_w, Inches(5.45))
    L = lambda lab: (lab + "：", {"sz": 13, "b": True, "c": s._accent})
    sat_label = "專案關係" if r["bucket"] == "C" else "滿足方式"
    first = True
    if r.get("asks"):
        para(body, [L("需求"), (r["asks"], {"sz": 13})], first=first, space_after=5, line=1.06); first=False
    para(body, [L(sat_label), (r["sat"], {"sz": 13})], first=first, space_after=5, line=1.06); first=False
    para(body, [L("資料來源"), (r["src"], {"sz": 13})], space_after=2, line=1.06)
    para(body, [("類別　", {"sz": 11.5, "b": True, "c": SOFT})] + class_runs(r["classes"]), space_after=5, line=1.05)
    para(body, [L("邊界"), (r["bnd"], {"sz": 13, "c": RGBColor(0x7A,0x2E,0x2E)})], space_after=5, line=1.06)
    if r.get("note"):
        para(body, [("註　", {"sz": 11.5, "b": True, "c": SOFT}), (r["note"], {"sz": 11.5, "i": True, "c": SOFT})], space_after=0, line=1.06)
    if has_img:
        if r.get("chart"):
            pic_framed(s, r["chart"], Inches(8.05), Inches(1.55), Inches(4.85), caption=r.get("cap"))
        else:
            pic_framed(s, r["img"], Inches(8.05), Inches(1.7), Inches(4.85), caption=r.get("cap"))
    footer(s, idx, total)
    return s

def slide_two_col(prs, idx, total, accent, title_runs, left_blocks, right_blocks, subtitle=None):
    s = add_blank(prs); bg(s, PAPER); s._accent = accent
    header(s, title_runs, subtitle)
    lt = tb(s, Inches(0.5), Inches(1.42), Inches(6.2), Inches(5.4))
    for i, blk in enumerate(left_blocks): para(lt, blk["runs"], first=(i==0), space_after=blk.get("sa",6), line=blk.get("ln",1.1), level=blk.get("lvl",0))
    rt = tb(s, Inches(6.95), Inches(1.42), Inches(5.95), Inches(5.4))
    for i, blk in enumerate(right_blocks): para(rt, blk["runs"], first=(i==0), space_after=blk.get("sa",6), line=blk.get("ln",1.1), level=blk.get("lvl",0))
    footer(s, idx, total)
    return s

def style_cell(cell, text, sz=10, bold=False, color=INK, fill=None, align=PP_ALIGN.LEFT, ea=EA_FONT):
    cell.margin_left = cell.margin_right = Inches(0.06)
    cell.margin_top = cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for ex in list(p.runs): ex.text = ""
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color
    _set_font(r, ea=ea)

def slide_table(prs, idx, total, accent, title_runs, headers, rows, col_w, subtitle=None, body_sz=9.5):
    s = add_blank(prs); bg(s, PAPER); s._accent = accent
    header(s, title_runs, subtitle)
    nrows, ncols = len(rows)+1, len(headers)
    top = Inches(1.45); left = Inches(0.45)
    width = sum(col_w); height = Inches(min(5.35, 0.34*nrows))
    gtbl = s.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = gtbl.table
    tbl.first_row = False; tbl.horz_banding = False
    # disable default style banding by setting our own fills
    for j, w in enumerate(col_w): tbl.columns[j].width = w
    for j, h in enumerate(headers):
        style_cell(tbl.cell(0, j), h, sz=10, bold=True, color=WHITE, fill=accent, align=PP_ALIGN.LEFT)
    for i, row in enumerate(rows, start=1):
        band = RGBColor(0xF4, 0xF6, 0xF8) if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            col = INK; bold=False; cellfill = band; al=PP_ALIGN.LEFT
            if isinstance(val, tuple):
                val, meta = val
                col = meta.get("c", INK); bold = meta.get("b", False)
                if meta.get("fill"): cellfill = meta["fill"]
                al = meta.get("al", PP_ALIGN.LEFT)
            style_cell(tbl.cell(i, j), val, sz=body_sz, bold=bold, color=col, fill=cellfill, align=al)
    footer(s, idx, total)
    return s

def slide_closing(prs, idx, total):
    s = add_blank(prs); bg(s, DARKBG)
    rect(s, Inches(0.9), Inches(1.35), Inches(0.9), Inches(0.10), BRAND)
    t = tb(s, Inches(0.9), Inches(0.75), Inches(11.5), Inches(0.7))
    para(t, [("誠實小結", {"sz": 30, "b": True, "c": WHITE})], first=True, space_after=0)
    body = tb(s, Inches(0.9), Inches(1.75), Inches(11.5), Inches(4.6))
    rows = [
        ("功能展示 ≠ 來源權威。", "多數 UI / 模型需求已在 demo route 功能展示；封包 / RF / 天氣 / 硬體證據仍受邊界限制或缺口。"),
        ("Tier A 營運商驗證 = 0 交付。", "全部 Bucket B 列均待 Tier A 換入；目前以 src/runtime/link-budget/ 之 Tier B 公開標準替代。"),
        ("本簡報絕不宣稱：", "封包延遲/抖動/吞吐/遺失 · 原生 RF 換手 / 主動服務路徑 · 營運商站台 RF 硬體真值 · 實測當地天氣 / R0.01 校準 · 外部驗收門檻 · 最終書面報告。"),
        ("可宣稱：", "公開來源可追溯 · 標準模型依揭露公式 · 幾何/可見性/模型選出換手候選 · 假設參數明確標示 · 24h soak 保存產物。"),
    ]
    for i,(h,b) in enumerate(rows):
        para(body, [(h+" ", {"sz": 15, "b": True, "c": RGBColor(0xF1,0xB0,0x70) if i==2 else BRANDLT}),
                    (b, {"sz": 13, "c": RGBColor(0xCF,0xD6,0xE0)})], first=(i==0), space_after=12, line=1.12)
    d = tb(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5))
    para(d, [("逐列來源 / 引用 / 邊界見 docs/requirement-presentation-walkthrough.md · 缺口登記見 docs/data-source-index.md",
              {"sz": 11, "c": RGBColor(0x8A,0x93,0xA6)})], first=True, space_after=0)
    footer(s, idx, total, dark=True)
    return s

# ============================================================ content (34 rows)
REQS_A = [
 dict(id="R1-T1 / K-A1", name="整合 LEO/MEO/GEO 軌道模型與多軌切換", bucket="A", status="complete",
   asks="整合三種軌道類別並產生多軌訊號切換。",
   sat="SGP4 傳播三軌道類別（manifest.json）；runtime-projection.ts 切換候選餵入 handover-policy.ts；demo route 可見。",
   src="1258 actors（651 OneWeb LEO + 33 Galileo MEO + 574 GEO 群組），CelesTrak 2026-06-13 快照 · [1][4][13]",
   classes=["geometric"],
   bnd="模型選出之切換候選，非主動服務路徑遷移或營運商事件記錄。",
   chart=ch("chart_actors.png"), cap="來源 1258 個 actor（runtime 截頂 200/100/60；幾何 / 公開來源）"),
 dict(id="R1-T2 / K-D1", name="整合衛星軌道資料", bucket="A", status="complete",
   asks="將真實軌道資料帶入模型。",
   sat="保存之 CelesTrak GP/TLE manifest + 傳播 actor 列；報告 Sources / Raw-data 分頁可追溯。",
   src="manifest.json 軌道類別計數 + epoch 範圍 · [1][4][5][13]",
   classes=["public"],
   bnd="快照無法證明上游即時目錄狀態或營運商指派。",
   img=IMG_DRAWER, cap="報告 Sources 分頁：軌道來源追溯"),
 dict(id="R1-T3 / K-D2", name="等效視覺場景呈現（Blender 或等效）", bucket="A", status="complete",
   asks="視覺化場景（Cesium 等效於原文所指 Blender 工具）。",
   sat="Cesium 球體 + 選定配對端點標記、相機框構、保存 HTML 報告。",
   src="投影 payload 之呈現 · [1][2][4][12]",
   classes=["display"],
   bnd="顯示轉換不產生新的來源真實性。",
   img=IMG_GLOBE, cap="選定配對 demo route（Cesium 球體）"),
 dict(id="R1-T4 / K-D3", name="互動式 UI", bucket="A", status="complete",
   asks="提供互動控制介面。",
   sat="V4 側欄（選擇器、降雨滑桿、統計、揭露區）、清單挑選、hover tooltip、copy-link；HTML 報告分頁 + CSV 下載。",
   src="app 介面互動 · [1][4][12]",
   classes=["display"],
   bnd="保存報告非完整 UI 互動軌跡。",
   img=IMG_REPORT, cap="V4 側欄 + 來源報告面板"),
 dict(id="R1-T5 / K-D4", name="換手規則與可調參數", bucket="A", status="complete",
   asks="設計換手規則與其驅動參數。",
   sat="handover-policy.ts（三策略 + 可調 gate），預設 cross-orbit-live；報告 Handover 分頁揭露門檻與事件。",
   src="引用 TR 38.821 §7.3.2.2 之策略 gate · [1][4][5][7]",
   classes=["standards"],
   bnd="模型化控制，非營運商 SLA 或封包軌跡。",
   img=IMG_HANDOVER, cap="模型選出之換手 / 跨軌遷移"),
 dict(id="R1-T6 / K-D5", name="通訊速率視覺化設計", bucket="A", status="complete",
   asks="視覺化通訊速率 / 連線品質。",
   sat="link-budget（FSPL + gas + rain + antenna）→ runtime-projection.ts 吞吐估計；側欄統計 + CSV。",
   src="晴空參考容量 + 大氣衰減 + 假設天線離軸折減 · [2][4][5][7][16]",
   classes=["standards","assumed"],
   bnd="延遲 / 抖動 / 吞吐為模型估計；營運商站台 RF 硬體真值（EIRP、天線盤、極化）仍為缺口。"),
 dict(id="R1-F1 / K-E1", name="支援 ≥500 LEO 模擬 actor（Starlink、OneWeb）", bucket="A", status="complete",
   asks="模擬至少 500 顆 LEO 衛星。",
   sat="預設載入 651 OneWeb LEO（+33 Galileo MEO、574 GEO 群組），CelesTrak 2026-06-13 快照（GROUP=oneweb|galileo|geo）。",
   src="651 OneWeb LEO ≥ 500 目標 · [1][4][5][13]",
   classes=["public","geometric"],
   bnd="「≥500 LEO」由 OneWeb 滿足；執行期 LEO 集為 OneWeb 非 Starlink。",
   note="「~11015」為上游 CelesTrak 目錄量，非本地載入（1258）— 2026-06-14 refresh（MEO 由 gnss 改回 galileo、GEO 為完整 geo 群組）。Starlink fixture 存在但非預設載入。",
   chart=ch("chart_buckets.png"), cap="需求分布：A 19 / B 7 / C 8（共 34）"),
 dict(id="R1-F2 / K-E2", name="可調模擬速度（即時 vs 預錄 TLE）", bucket="A", status="complete",
   asks="可調播放速度與即時 / 預錄場景切換。",
   sat="HUD 三段預設 30× / 60× / 120×（m8a-v4-product-ux-model.ts）；即時視窗預設 [now, now+360 min]；demo route 為預錄固定視窗。",
   src="預設常數 + 6 小時預錄視窗 · [1][4][5]",
   classes=["display"],
   bnd="播放預設為 app 介面證據，非列級 CSV 指標；底層時鐘接受任意倍率，可見 UI 限三段。"),
 dict(id="R1-F3 / K-E3", name="即時可通訊時間（原文提及 iperf/ping）", bucket="A", status="partial",
   asks="顯示即時可通訊時間；原文提及 iperf/ping 測試。",
   sat="runtime-projection.ts 可見視窗 + 連線選擇策略；報告 Summary 顯示、CSV 匯出。",
   src="模型計算之可用時間 · [1][2][4][5]",
   classes=["geometric","standards"],
   bnd="為模型計算之可用性，非成功連線時間 / ping / iperf uptime / 營運商可用率。iperf/ping 量測子項為資料缺口（irreducible-1）。",
   img=IMG_VIS, cap="可見視窗 → 模型可通訊時間"),
 dict(id="R1-F4 / K-E4 / K-F4", name="依延遲/抖動/速率切換換手策略", bucket="A", status="complete",
   asks="依延遲 / 抖動 / 網速條件在高/中/低軌間切換換手策略。",
   sat="link-budget 指標餵入換手引擎；報告 Handover 分頁之模型選出切換 + 跨軌遷移。",
   src="引用 TR 38.821 §7.3.2.2 + V-MO1 之策略 gate · [1][4][5][7]",
   classes=["standards"],
   bnd="無營運商事件記錄、主動服務路徑或 RF 換手軌跡。"),
 dict(id="R1-F5 / K-E5", name="統計報告匯出", bucket="A", status="complete",
   asks="匯出統計。",
   sat="V4 側欄 CSV（RFC-4180，runtime-projection-csv.ts）+ 獨立 HTML 報告（runtime-projection-evidence-report.ts）；皆含 SHA-256 保存於 [3]。",
   src="報告 HTML（1646168 B, sha 05bf3293…）+ CSV（104776 B, sha 0c133196…）· [1][3][4][5]",
   classes=["display","external"],
   bnd="HTML / CSV 為證據面，非最終驗收證明。",
   img=IMG_REPORT, cap="HTML 報告 + 列級 CSV（含 SHA-256）"),
 dict(id="K-A4", name="輸入 TLE 並追蹤衛星運動", bucket="A", status="complete",
   asks="輸入 TLE 資料並追蹤衛星。",
   sat="TLE/OMM manifest 列傳播入 actor provenance；chrome TLE 晶片由最新保存 TLE 時戳推導（resolveTleDate, chrome-telemetry.ts）。",
   src="晶片顯示 2026-05-31；CelesTrak 交叉核對 NORAD 44057 = ONEWEB-0012 · [1][4][5][13]",
   classes=["public"],
   bnd="瀏覽器讀取打包檔，非證明即時目錄連線；無營運商驗證主張。",
   note="原 G1 唯一失敗列；merge Task 1 已修。"),
 dict(id="K-E6", name="呈現降雨衰減影響", bucket="A", status="complete",
   asks="展示降雨衰減影響。",
   sat="rain-attenuation.ts（完整 P.618-14 path method）+ V4 降雨滑桿；下行吞吐隨雨折減（高雨率例：LEO −45% / GEO −80%）。",
   src="ITU-R P.618-14 §2.2.1.1，例 20 GHz / 50 mm/h / 45° → 27.3 dB；隨緯度（0°→30.5、78°→8.3）· [1][2][4][5][7]",
   classes=["standards"],
   bnd="由互動瞬時雨率（what-if 輸入）驅動，非實測當地天氣或校準雨率時間序列。",
   chart=ch("chart_rain.png"), cap="P.618-14 模型錨點（範例輸入，非實測）"),
 dict(id="K-F7", name="產出 demo（預錄場景或即時模擬）", bucket="A", status="complete",
   asks="交付可執行 demo。",
   sat="選定配對 V4 demo route + 保存證據包。",
   src="route + 投影輸出 · [1][3][4][5]",
   classes=["display","external"],
   bnd="代表性 demo 證據非外部驗收門檻腳本。"),
 dict(id="R1-D1", name="11/30 里程碑：軌道模型匯入成功", bucket="A", status="complete",
   sat="匯入軌道來源檔 + 傳播 actor，報告揭露。",
   src="manifest + actor 列 · [1][4][5][13]",
   classes=["public","geometric"],
   bnd="本列不自證完整歷史里程碑包。"),
 dict(id="R1-D2", name="11/30 里程碑：動態參數 UI", bucket="A", status="complete",
   sat="V4 側欄保留時間視窗、降雨輸入、作用策略。",
   src="快照參數 · [1][4][5]",
   classes=["display"],
   bnd="快照參數非完整即時控制互動記錄。"),
 dict(id="R1-D3", name="11/30 里程碑：通訊時間統計", bucket="A", status="complete",
   sat="報告顯示 + CSV 匯出通訊時間統計。",
   src="communicationStats payload · [1][4][5]",
   classes=["geometric"],
   bnd="統計為模型推導，非 ping / iperf 或實測吞吐。"),
 dict(id="R1-D4", name="11/30 里程碑：穩定運行 ≥24 小時", bucket="A", status="complete",
   sat="保存之 24 小時 soak 產物。",
   src="output/soak/2026-05-15…/summary.json — passed=true, durationMs=86400000(24h), sampleCount=1289, failureCount=0 · [1][4]",
   classes=["external"],
   bnd="開啟報告不會重跑 soak；主張依賴保存產物。"),
 dict(id="V-MO1", name="跨軌 LIVE 換手（K-A1 口頭補充）", bucket="A", status="complete",
   asks="單一服務連續換手 LEO → MEO → GEO（真實 live 換手，非軌道類型選擇器）。",
   sat="handover-policy.ts cross-orbit-live 為 demo route 預設且唯一策略（opt-in 旗標已移除）；跨軌遷移時 replay pill 帶 data-v-mo1=\"true\"。",
   src="模型選出之跨軌遷移事件 · [1][4][5][7]",
   classes=["geometric","standards"],
   bnd="不可稱原生 RF 換手、主動服務路徑遷移或營運商事件證明；為 TLE 幾何上之模型策略遷移。",
   img=IMG_LINK, cap="跨軌模型策略遷移（V-MO1）"),
]

REQS_B = [
 dict(id="K-A2", name="連線品質規則（延遲/抖動/速率切換策略）", bucket="B", status="partial",
   asks="營運商連線品質規則驅動軌道切換。",
   sat="runtime-projection.ts 延遲模型 + runtime-modeled-output.ts 抖動模型 + handover-policy.ts（Tier B 替代）。",
   src="單向延遲 TR 38.811 §6.6.2（球面斜距 Eq 6.6-3）+ clause 5.3.1.1；抖動 = 軌道基線 × 雨率 · [1][2][4][5][7]",
   classes=["standards"],
   bnd="無營運商封包軌跡或私有連線品質序列；延遲為模型錨點，非 ping / RTT / jitter 量測。+2 ms 處理項揭露為非標準。",
   note="延遲引用更正為 §6.6.2 + clause 5.3.1.1（§6.7 為 fast-fading 非延遲）。"),
 dict(id="K-A3-a", name="天線參數（峰值增益、波束寬、場型）", bucket="B", status="partial",
   asks="天線峰值增益、波束寬與輻射場型。",
   sat="antenna-pattern.ts（computeSatellite/EarthStationAntennaGainDb）已接入 runtime-projection.ts（line 610/616），經 runtime-antenna-assumptions.ts。",
   src="ITU-R S.1528-0 Annex 1 + S.465-6，假設每軌參數（LEO 35/43、MEO 38/48、GEO 42/52 dBi）· [1][2][4][5][7][16]",
   classes=["standards","assumed"],
   bnd="RSRP 匯出為 receivedPowerProxyDbm（相對 proxy，EIRP 未知）；假設增益/波束/盤徑非營運商 RF 硬體、真實盤徑、EIRP、極化、波束指派或實測天線增益。",
   note="舊 inventory/spine 稱「standalone 未接入」為過時；現已 runtime-wired，剩餘缺口為 operator-rf-hardware-truth。",
   chart=ch("chart_antenna.png"), cap="合併增益 75.7 / 83.9 / 91.9 dB（模型 + 假設）"),
 dict(id="K-A3-b", name="降雨衰減模型", bucket="B", status="partial",
   asks="降雨衰減模型。",
   sat="rain-attenuation.ts（完整 P.618-14 path method）；P.838-3 係數委派 itu-r-p838-rain-attenuation.ts（係數表單一真源）。",
   src="P.618-14 §2.2.1.1：γR=k·R^α(S5)+ r0.01(Eq5)+ ν0.01(Eq6,用站緯)+ LE=LR·ν0.01(S8)+ A=γR·LE(S9)；P.838-3 對 Table 5 驗證一致 · [1][2][4][5][7]",
   classes=["standards"],
   bnd="由瞬時雨率輸入驅動，非 R0.01 長期統計；無時間百分比外推（S10）；P.839 雨高簡化。非實測當地天氣或台灣在地校準。"),
 dict(id="K-F2", name="整合 V-team 模擬程式（天線 + 降雨 + ITU 規則）", bucket="B", status="partial",
   asks="整合夥伴 V-team 模擬元件。",
   sat="自寫 ITU 計算器跨五個 link-budget 模組（FSPL + rain + gas + antenna + handover-policy），均已接入 runtime-projection.ts。",
   src="P.618 / P.676 / S.1528 / S.465 / TR 38.811 / TR 38.821 · [1][2][4][5][7][16]",
   classes=["standards","assumed"],
   bnd="為本專案對 V-team 程式之替代，非 V-team 程式本身；標準無法補足缺少之站台硬體真值。待 V-team / 營運商資料交付後換為 Tier A。"),
 dict(id="irreducible-1", name="真實封包軌跡（延遲/抖動時間序列）", bucket="B", status="gap",
   asks="真實 ESTNeT 封包軌跡與延遲/抖動時間序列。",
   sat="僅 TR 38.811 generic model 之合成 / 模型基線；truth-boundary 標 \"ESTNeT trace pending\"。",
   src="模型基線 · [1][2][4][5][6]",
   classes=["standards","gap"],
   bnd="無保存 ping / iperf 或等效封包測試時間序列；不得宣稱實測延遲 / 抖動。"),
 dict(id="irreducible-2", name="外部驗收場景腳本 + pass/fail 門檻", bucket="B", status="gap",
   asks="驗收測試腳本（哪個案例通過、門檻為何）。",
   sat="五條代表性選定配對基線 route 保存於 v4-projection-wave1-baselines.ts（支持 demo 案例討論）。",
   src="保存基線 · [1][2][4][6]",
   classes=["external","gap"],
   bnd="代表性 route 非外部驗收門檻腳本、場景包或審查裁定。"),
 dict(id="irreducible-3", name="台灣在地雨量統計校準（R0.01）", bucket="B", status="partial",
   asks="在地台灣雨量統計（如 R0.01）供校準。",
   sat="保存之公開 CWA 站觀測樣本映射至有界 demo 預設。",
   src="CWA 46691 Anpu + 46693 Yangmingshan，各 149 樣本，2026-06-12→13，raw curl + sha ×2；峰值 10-min 0.5 mm ×6=3 mm/h → 5 mm/h 預設 · [1][2][6][11][14]",
   classes=["public","assumed"],
   bnd=f"非 CHT-SANSA 連線實測天氣、非 SANSA 降雨、非 {ds.DATE_YMD} 場景視窗天氣、非 R0.01 / 可用率校準；晴空報告雨率仍為 0 除非移動滑桿。",
   note="舊 inventory 標純缺口為過時；公開樣本已存在，剩 measured-for-link / R0.01 強度為缺口。"),
]

REQS_C = [
 dict(id="K-F1", name="經外部模擬器之實體網路橋接（INET↔ESTNeT）", bucket="C", status="scope",
   sat="外部基礎設施；無選定配對 runtime 內容。",
   src="模型可見性 / 換手無法證明實體橋接 · [1][2][4]",
   classes=["gap"], bnd="（外部）模型可見性 / 換手無法證明實體橋接。"),
 dict(id="K-A5", name="Linux 主環境 + Windows + WSL 備援", bucket="C", status="scope",
   sat="屬外部驗證，非報告層。",
   src="非 UI / runtime 投影值 · [1][4]",
   classes=["external"], bnd="非 UI / runtime 投影值；屬主機環境就緒。"),
 dict(id="K-B1", name="Windows tunneling / ESTNeT 真流量衛星橋接", bucket="C", status="scope",
   asks="ping 10.2.0.1 → tun0 → GS-A → Satellite → GS-B → tun1 拓樸。",
   sat="拓樸圖，非 UI 規格。",
   src="無 tunnel / 介面 / 真流量橋接證據 · [1][2][4]",
   classes=["gap"], bnd="（外部）無 tunnel / 介面 / 真流量橋接證據附上。"),
 dict(id="K-C1", name="INET NAT 路由 + veth 橋接", bucket="C", status="scope",
   sat="外部拓樸產物。",
   src="原始投影 JSON 無法證明 NAT 路由 · [1][2][4]",
   classes=["gap"], bnd="（外部）原始投影 JSON 無法證明 NAT 路由。"),
 dict(id="K-F5", name="虛擬 DUT testbench 程式", bucket="C", status="scope",
   sat="非選定配對包之一部分。",
   src="投影資料非 testbench 程式 · [1][2][4]",
   classes=["gap"], bnd="（外部）投影資料非 testbench 程式。"),
 dict(id="K-F6", name="實體 DUT / NE-ONE 流量產生器場景", bucket="C", status="scope",
   sat="無實體 DUT 或流量結果保存。",
   src="不得宣稱實測吞吐或流量產生成功 · [1][2][4]",
   classes=["gap"], bnd="（外部）不得宣稱實測吞吐或流量產生器成功。"),
 dict(id="K-F8", name="最終書面報告", bucket="C", status="gap",
   sat="由報告層撰寫；選定配對 HTML 為附錄候選，非最終報告。",
   src="無最終書面報告包 · [1][2][3][4]",
   classes=["gap"], bnd="（報告層）無最終書面報告包附上。"),
 dict(id="R1-D5", name="技術 WP1 評估報告", bucket="C", status="gap",
   sat="選定配對報告可支持未來 WP1，但不能取代。",
   src="無完整技術評估報告產物 · [1][2][3][4]",
   classes=["gap"], bnd="（報告層）無完整技術評估報告產物附上。"),
]

GAP_ROWS = [
 ("封包延遲、抖動、吞吐、遺失", "R1-F3/K-E3；K-A2；irr-1", "保存 ping/iperf 記錄：端點、方向、時戳、raw、解析指標、審查裁定。"),
 ("原生 RF 換手 / 主動服務路徑", "R1-F4/K-E4/K-F4；V-MO1", "營運商 / RF 事件記錄、終端記錄、實驗室 RF 軌跡或主動路徑遙測。"),
 ("營運商站台 RF 硬體", "R1-T6/K-D5；K-A3-a；K-F2", "列級 EIRP、增益、波束寬、極化、頻段、天線型號（operator-rf-hardware-truth）。"),
 ("勘測 RF 視界 / 全表 DEM", "站台精度；K-F2 脈絡", "勘測視界資料，或將已審 DEM 取樣演算法擴及選定配對以外（現僅 CHT+SANSA）。"),
 ("實測連線天氣 + 長期 R0.01", "K-A3-b；K-E6；irr-3", "長期站/網格統計或場景視窗天氣；R0.01 可用率校準。"),
 ("外部驗收門檻 + 最終驗收", "irr-2；K-F8；R1-D5", "驗收場景包：案例清單、門檻、指令、raw、pass/fail；正式書面報告。"),
 ("實體橋接 / NAT / DUT / 流量產生器", "K-F1；K-B1；K-C1；K-F5；K-F6", "拓樸清單、介面/NAT/路由設定、橋接/tunnel 記錄、DUT 識別、流量輸出、遮蔽策略、審查裁定。"),
 ("Tier A 營運商資料（Bucket B 換入）", "K-A2；K-A3-a；K-A3-b；K-F2", "營運商提供之保存資料 + 驗證記錄 + Tier B→A 換入政策。"),
]

CATALOG = [
 ("[1]", "requirements-consolidated.md", "34 列清單、bucket 計數、ID、tier 政策。", "runtime 證明、封包/RF 量測、驗收。"),
 ("[2]", "docs/data-source-index.md", "來源分類、呈現來源圖、缺口修復、「模型仍是模型」。", "原始證據庫、驗收證書。(對天線/標準節為現行權威)"),
 ("[3]", "selected-pair-source-evidence/README.md", "保存選定配對包、route、產生時間、檔案雜湊、邊界。", "封包/RF/營運商驗證、最終書面驗收。"),
 ("[4]", "runtime-projection-evidence-…-360m.html", "可讀報告：需求、來源、模型、換手、可見性、raw。", "原始證據庫、重跑證明、封包/RF 量測、正式驗收。"),
 ("[5]", "runtime-projection-…-360m.csv", "列級匯出：欄位、來源 tier、站台精度、模型輸出、缺口、非主張。", "payload 以外敘事、封包記錄、驗收。"),
 ("[6]", "external-source-reconciliation.md", "CHT/SANSA 公開來源對帳、衝突、裁定、修復路徑、不可觀測缺口。", "CHT 精確座標、主動路徑、原生 RF 換手、封包量測、驗收。"),
 ("[7]", "3gpp-itu-references/README.md + PDFs", "3GPP/ITU-R PDF 對照、節引用、checksum。", "營運商驗證、實測品質、當地天氣、硬體、封包/RF 結果。"),
 ("[8]", "multi-orbit-public-registry-sources.md", "公開站台能力、來源 URL、座標/高程/頻段、非主張。", "商業配對路由、主動服務衛星、私有設施幾何、RF 硬體完整性。"),
 ("[9]", "open-elevation-…-2026-06-12.json", "重現選定配對舊高程快取值（Open-Elevation）。", "官方 DEM tile/cell/datum/版本/授權；地形遮罩；CHT 精確座標。"),
 ("[10]", "copernicus-dem-…-sample-2026-06-12.json", "選定配對 Copernicus GLO-30 比對 cell（tile/row/col/CRS/datum/license）。", "全表替換；封包/RF/路由/驗收。"),
 ("[11]", "rain-source-repair-candidates-2026-06-12.json", "候選公開雨源（NASA GPM IMERG、CWA Anbu/Zhuzihu）。", "任何取樣雨值、實測天氣、雨衰驗證。"),
 ("[12]", "README.md", "repo 範圍、資料邊界規則、指令清單。", "需求覆蓋細節或原始證明。"),
 ("[13]", "satellites-network/manifest.json", "CelesTrak 軌道 manifest、產生時間、軌道類別計數、epoch 範圍。", "即時上游目錄、主動服務衛星、營運商指派。"),
 ("[14]", "local-rain-calibration-2026-06-13.json", "公開 CWA 統計（46691 Anpu、46693 Yangmingshan），raw curl + sha ×2，各 149 樣本 → 5 mm/h 預設。", "實測連線天氣、SANSA 降雨、場景視窗天氣、R0.01 校準、雨衰驗證。"),
 ("[15]", "copernicus-dem-…-terrain-mask-2026-06-13.json", "Copernicus GLO-30 高程採用（CHT 489 m）+ 地形遮罩 radial-nearest-cell-horizon-v1 → CHT 21° / SANSA 1°。", "勘測 RF 視界、遮蔽物、營運商實測高程、全 69 列替換。"),
 ("[16]", "runtime-antenna-assumptions.ts", "揭露之 assumed-tier-b-antenna-params-selected-pair-v1，接入選定配對 RSRP / 吞吐。", "營運商盤徑、EIRP、極化、波束寬、天線型號、實測增益。"),
]

# ============================================================ assemble
def cls_short(keys):
    return "／".join(CLS[k][0] for k in keys)

def coverage_rows(reqs):
    rows = []
    for r in reqs:
        color, _ = ST[r["status"]]
        st_zh = {"complete":"完成","partial":"部分(B)","gap":"缺口","scope":"範圍外"}[r["status"]]
        rows.append([r["id"], r["name"], (st_zh, {"c": color, "b": True}), cls_short(r["classes"])])
    return rows

def main():
    build_charts()
    prs = Presentation(TEMPLATE)
    # strip sample slides (keep masters/layouts/theme): drop sldId AND its
    # relationship+part, else orphaned slide parts collide with new slideN.xml.
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        rId = sid.get(qn("r:id"))
        if rId:
            try: prs.part.drop_rel(rId)
            except Exception: pass
        lst.remove(sid)

    deck = []   # list of builder callables(idx,total)
    # front matter
    deck.append(lambda i,t: slide_title(prs))
    deck.append(lambda i,t: slide_two_col(prs, i, t, BRAND,
        [("大綱 ", {"sz": 26, "b": True, "c": INK}), ("Agenda", {"sz": 18, "b": True, "c": BRAND})],
        [ {"runs":[("1　誠實閱讀規則與來源分級", {"sz":15,"b":True,"c":INK})], "sa":10},
          {"runs":[("2　需求總覽（34 列：A 19 / B 7 / C 8）", {"sz":15,"b":True,"c":INK})], "sa":10},
          {"runs":[("3　選定配對 demo route", {"sz":15,"b":True,"c":INK})], "sa":10},
          {"runs":[("4　Bucket A：專案 demo 19 列", {"sz":15,"b":True,"c":INK})], "sa":10},
          {"runs":[("5　Bucket B：Tier B 公開標準替代 7 列", {"sz":15,"b":True,"c":INK})], "sa":10},
        ],
        [ {"runs":[("6　Bucket C：外部 / 報告層 8 列", {"sz":15,"b":True,"c":INK})], "sa":10},
          {"runs":[("7　資料缺口總表（已揭露待補）", {"sz":15,"b":True,"c":INK})], "sa":10},
          {"runs":[("8　完整 34 列對應表（附錄）", {"sz":15,"b":True,"c":INK})], "sa":10},
          {"runs":[("9　來源目錄 [1]–[16]（附錄）", {"sz":15,"b":True,"c":INK})], "sa":10},
          {"runs":[("10　誠實小結", {"sz":15,"b":True,"c":INK})], "sa":10},
        ]))
    # honesty rule + source classes
    def rule_slide(i,t):
        s = add_blank(prs); bg(s, PAPER); s._accent = BRAND
        header(s, [("誠實閱讀規則 ", {"sz": 26, "b": True, "c": INK}), ("每個數字問兩個問題", {"sz": 16, "b": True, "c": BRAND})])
        lt = tb(s, Inches(0.5), Inches(1.45), Inches(5.6), Inches(5.3))
        para(lt, [("1　顯示值從何而來？", {"sz": 15, "b": True, "c": INK})], first=True, space_after=4)
        para(lt, [("（公開產物 / 標準模型 / TLE 幾何 / 缺口）", {"sz": 12, "c": SOFT})], space_after=12)
        para(lt, [("2　該來源實際能支撐什麼主張？", {"sz": 15, "b": True, "c": INK})], space_after=12)
        para(lt, [("功能展示 ≠ 來源權威。", {"sz": 14, "b": True, "c": RGBColor(0xC6,0x28,0x28)})], space_after=6, line=1.15)
        para(lt, [("多數需求已由 viewer 功能展示，但其「實測 / 營運商驗證」之較強用語仍為缺口。簡報絕不將模型 / 假設 / 公開推導值呈現為實測或營運商驗證。",
                   {"sz": 12.5, "c": INK})], space_after=12, line=1.2)
        para(lt, [("來源權威分級：", {"sz": 13, "b": True, "c": INK})], space_after=4)
        para(lt, [("Tier A 營運商驗證 = 0 交付　·　Tier B 公開標準（src/runtime/link-budget/）　·　Tier C 幾何（SGP4 + TLE）",
                   {"sz": 12, "c": SOFT})], space_after=0, line=1.15)
        # right: 8 source-class chips
        rt = tb(s, Inches(6.45), Inches(1.45), Inches(6.4), Inches(0.4))
        para(rt, [("資料來源分類（8 類）", {"sz": 13, "b": True, "c": INK})], first=True, space_after=0)
        y = 1.92
        for k in ["public","standards","geometric","assumed","display","external","legacy","gap"]:
            zh,en,col = CLS[k]
            rect(s, Inches(6.45), Inches(y), Inches(0.22), Inches(0.5), col)
            ct = tb(s, Inches(6.78), Inches(y-0.02), Inches(6.0), Inches(0.56), anchor=MSO_ANCHOR.MIDDLE)
            para(ct, [(f"{zh}  ", {"sz": 12.5, "b": True, "c": col}), (en, {"sz": 10.5, "b": True, "c": col})], first=True, space_after=0)
            y += 0.62
        footer(s, i, t)
        return s
    deck.append(rule_slide)
    # overview
    def overview(i,t):
        s = add_blank(prs); bg(s, PAPER); s._accent = BRAND
        header(s, [("需求總覽 ", {"sz": 26, "b": True, "c": INK}), ("34 列 · A 19 / B 7 / C 8", {"sz": 16, "b": True, "c": BRAND})])
        lt = tb(s, Inches(0.5), Inches(1.45), Inches(6.4), Inches(5.3))
        para(lt, [("A　專案 demo（19）", {"sz": 15, "b": True, "c": BRAND})], first=True, space_after=3)
        para(lt, [("Cesium 球體 + 選定配對 V4 投影；幾何 / 標準推導，顯示轉換。無主動服務路徑、封包量測或營運商遙測主張。", {"sz": 12, "c": INK})], space_after=10, line=1.18)
        para(lt, [("B　Tier B 公開標準替代（7）", {"sz": 15, "b": True, "c": RGBColor(0xC9,0x6A,0x00)})], space_after=3)
        para(lt, [("各列需營運商資料達 Tier A；以 ITU-R / 3GPP 公開標準自實作替代（4 已 runtime-wired，3 irreducible）。", {"sz": 12, "c": INK})], space_after=10, line=1.18)
        para(lt, [("C　外部基礎設施 / 報告層（8）", {"sz": 15, "b": True, "c": RGBColor(0x6B,0x72,0x80)})], space_after=3)
        para(lt, [("不在選定配對 viewer 範圍：6 外部網路/測試基礎設施 + 2 報告層交付。列出以求需求清單完整。", {"sz": 12, "c": INK})], space_after=0, line=1.18)
        pic_framed(s, ch("chart_buckets.png"), Inches(7.2), Inches(1.7), Inches(5.7), caption="需求分布（共 34 列）")
        bigt = tb(s, Inches(7.2), Inches(5.2), Inches(5.7), Inches(1.3))
        para(bigt, [("34", {"sz": 54, "b": True, "c": BRAND}), ("　需求　·　", {"sz": 16, "c": INK}),
                    ("0", {"sz": 54, "b": True, "c": RGBColor(0xC6,0x28,0x28)}), ("　Tier A 交付", {"sz": 16, "c": INK})], first=True, space_after=0)
        footer(s, i, t)
        return s
    deck.append(overview)
    # demo route
    def demo(i,t):
        s = add_blank(prs); bg(s, PAPER); s._accent = BRAND
        header(s, [("選定配對 demo route ", {"sz": 24, "b": True, "c": INK}), ("CHT Yangmingshan ↔ SANSA Hartebeesthoek", {"sz": 15, "b": True, "c": BRAND})])
        lt = tb(s, Inches(0.5), Inches(1.45), Inches(5.4), Inches(5.3))
        para(lt, [("路由：", {"sz": 13, "b": True, "c": INK})], first=True, space_after=3)
        para(lt, [(ds.DEMO_ROUTE,
                   {"sz": 11, "c": BRAND, "ea": LAT_FONT})], space_after=10, line=1.2)
        para(lt, [("視窗：", {"sz": 13, "b": True, "c": INK}), (ds.WINDOW_LABEL_ZH, {"sz": 12.5, "c": INK})], space_after=8)
        para(lt, [("來源 tier：", {"sz": 13, "b": True, "c": INK}), ("geometric-derived；evidence kind = cross-family-geometric。", {"sz": 12.5, "c": INK})], space_after=8, line=1.15)
        para(lt, [("標籤掃描：", {"sz": 13, "b": True, "c": INK}), ("HTML 報告 39 standards-derived + 12 geometric-derived chip，0 operator-validated / public-disclosed。", {"sz": 12.5, "c": INK})], space_after=8, line=1.18)
        para(lt, [("即：產生器不將任何模型 / 幾何值升級為營運商權威。", {"sz": 12.5, "b": True, "c": RGBColor(0x2E,0x7D,0x32)})], space_after=0, line=1.15)
        pic_framed(s, IMG_GLOBE, Inches(6.2), Inches(1.55), Inches(6.7), caption="Cesium 球體 · 選定配對端點 · 模型換手（retained screenshot）")
        footer(s, i, t)
        return s
    deck.append(demo)
    # Section A
    deck.append(lambda i,t: slide_section(prs, i, t, "Bucket A · 19 列", "專案 demo 列",
        ["主軸：Cesium 球體 viewer + 選定配對 V4 地面站投影。",
         "誠實來源類別：幾何推導（TLE + SGP4 + 可見性）+ 標準模型推導（link budget），經顯示轉換（面板、報告、CSV）。",
         "本區無任何列宣稱主動服務路徑、封包量測或營運商遙測。"], BRAND))
    for r in REQS_A: deck.append(lambda i,t,r=r: slide_req(prs, i, t, r))
    # Section B
    deck.append(lambda i,t: slide_section(prs, i, t, "Bucket B · 7 列", "Tier B 公開標準替代",
        ["各列需營運商資料才能達 Tier A 權威；本專案於 src/runtime/link-budget/ 自實作 Tier B 公開標準替代，附保存之 ITU-R / 3GPP PDF + checksum [7]。",
         "4 替代已接入選定配對 runtime；3 條 irreducible 有合理公開替代但其最強用語仍為缺口。",
         "2026-06-13 merge 改變三列：天線已 runtime-wired、已有公開在地雨統計、選定配對已採 DEM / 地形 provenance。"], RGBColor(0xE0,0x8A,0x1E)))
    for r in REQS_B: deck.append(lambda i,t,r=r: slide_req(prs, i, t, r))
    # Section C
    deck.append(lambda i,t: slide_section(prs, i, t, "Bucket C · 8 列", "外部基礎設施 / 報告層",
        ["不在選定配對 viewer 之 UI 範圍。6 列為本專案以外擁有之外部網路 / 測試基礎設施；2 列為本專案另行撰寫之報告層交付。",
         "列出以求需求清單完整。選定配對投影為球體 / runtime 模型產物，無法自證以下任一列。"], RGBColor(0x6B,0x72,0x80)))
    for r in REQS_C: deck.append(lambda i,t,r=r: slide_req(prs, i, t, r))
    # gap table
    deck.append(lambda i,t: slide_table(prs, i, t, RGBColor(0xC6,0x28,0x28),
        [("資料缺口總表 ", {"sz": 24, "b": True, "c": INK}), ("需營運商 / 外部資料（已揭露待補，不在簡報中宣稱）", {"sz": 14, "b": True, "c": RGBColor(0xC6,0x28,0x28)})],
        ["缺口", "阻擋需求", "關閉所需"],
        [[ (g[0],{"b":True}), g[1], g[2]] for g in GAP_ROWS],
        [Inches(2.7), Inches(2.9), Inches(6.85)], body_sz=9.5))
    # coverage appendix (A) and (B+C)
    deck.append(lambda i,t: slide_table(prs, i, t, BRAND,
        [("附錄：完整對應表（1/2）", {"sz": 22, "b": True, "c": INK}), ("　Bucket A · 19 列", {"sz": 14, "b": True, "c": BRAND})],
        ["需求 ID", "需求", "狀態", "主來源類別"],
        coverage_rows(REQS_A),
        [Inches(2.5), Inches(6.0), Inches(1.4), Inches(2.55)], body_sz=8.6))
    deck.append(lambda i,t: slide_table(prs, i, t, RGBColor(0xE0,0x8A,0x1E),
        [("附錄：完整對應表（2/2）", {"sz": 22, "b": True, "c": INK}), ("　Bucket B · 7 + Bucket C · 8", {"sz": 14, "b": True, "c": RGBColor(0xE0,0x8A,0x1E)})],
        ["需求 ID", "需求", "狀態", "主來源類別"],
        coverage_rows(REQS_B + REQS_C),
        [Inches(2.5), Inches(6.0), Inches(1.4), Inches(2.55)], body_sz=8.6))
    # source catalog appendix (2 slides x 8)
    deck.append(lambda i,t: slide_table(prs, i, t, BRAND,
        [("附錄：來源目錄 [1]–[8]", {"sz": 22, "b": True, "c": INK})],
        ["ID", "來源", "可支持", "不可支持"],
        [[ (c[0],{"b":True,"c":BRAND}), c[1], c[2], (c[3],{"c":RGBColor(0x7A,0x2E,0x2E)})] for c in CATALOG[:8]],
        [Inches(0.7), Inches(3.5), Inches(4.55), Inches(3.7)], body_sz=8.4))
    deck.append(lambda i,t: slide_table(prs, i, t, BRAND,
        [("附錄：來源目錄 [9]–[16]", {"sz": 22, "b": True, "c": INK})],
        ["ID", "來源", "可支持", "不可支持"],
        [[ (c[0],{"b":True,"c":BRAND}), c[1], c[2], (c[3],{"c":RGBColor(0x7A,0x2E,0x2E)})] for c in CATALOG[8:]],
        [Inches(0.7), Inches(3.5), Inches(4.55), Inches(3.7)], body_sz=8.4))
    # closing
    deck.append(lambda i,t: slide_closing(prs, i, t))

    total = len(deck)
    for idx, fn in enumerate(deck, start=1):
        fn(idx, total)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"slides: {total}  ->  {OUT}")

if __name__ == "__main__":
    main()
