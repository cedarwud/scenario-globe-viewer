#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
WS-D — Evidence Legibility slide deck (standalone, build-only; not committed).

Four content slides (D1-D4) + a title slide that make the viewer's authenticity
work legible: prove live compute to a viewer, and prove auditability to a
procurement acceptance reviewer.

Honesty rule: every number here is modeled / standard-derived / a disclosed
proxy. Nothing is presented as measured / operator-validated / active-serving;
the external gaps are disclosed, not invented. All figures are sourced from the
regenerated retained evidence package (deliverable/selected-pair-source-evidence/)
and the runtime physics modules — see the per-figure provenance.

Self-contained: it does NOT import or modify the finalized requirement decks.

Run:  python3 scripts/build_evidence_legibility_slides.py
Out:  deliverable/evidence-legibility-slides.pptx
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import demo_scenario as ds   # demo window/route single source of truth

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "deliverable", "evidence-legibility-slides.pptx")
# D2 pipeline figure rendered by research-visual-lab (viz.json -> SVG -> PNG).
# If absent, slide D2 falls back to a hand-drawn box pipeline.
D2_FIG = os.path.join(ROOT, "deliverable", "evidence-legibility-figures", "d2-compute-pipeline.png")

# Palette (matches the requirement decks' clean navy/teal system).
INK    = RGBColor(0x1A, 0x22, 0x38)   # deep navy text
DARKBG = RGBColor(0x12, 0x19, 0x2B)   # title / closing bg
PAPER  = RGBColor(0xFF, 0xFF, 0xFF)   # content bg
SOFT   = RGBColor(0x5B, 0x63, 0x72)   # muted gray
HAIR   = RGBColor(0xD9, 0xDD, 0xE3)   # hairline
BRAND  = RGBColor(0x2A, 0x9D, 0x8F)   # teal accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

# Fidelity-tier colors (mirror the report's honesty grading).
REAL     = RGBColor(0x2E, 0x7D, 0x32)  # real compute (green)
STANDARD = RGBColor(0x1B, 0x80, 0x74)  # standard model (teal)
PROXY    = RGBColor(0xC9, 0x6A, 0x00)  # illustrative proxy (amber)
GAP      = RGBColor(0xC6, 0x28, 0x28)  # external gap (red)

EA_FONT = "Microsoft JhengHei"
LAT_FONT = "Arial"
MONO = "Consolas"

SW, SH = Inches(13.333), Inches(7.5)


def _set_font(run, latin=LAT_FONT, ea=EA_FONT):
    run.font.name = latin
    rpr = run._r.get_or_add_rPr()
    import copy
    for tag in ("a:latin", "a:ea", "a:cs"):
        existing = rpr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag.split(":")[1])
        if existing is not None:
            rpr.remove(existing)
    from pptx.oxml.ns import qn
    for tag, font in (("a:latin", latin), ("a:ea", ea), ("a:cs", latin)):
        el = rpr.makeelement(qn(tag), {qn("w:typeface") if False else "typeface": font})
        rpr.append(el)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    r = slide.shapes.add_shape(shape, l, t, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line
        r.line.width = Pt(0.75)
    r.shadow.inherit = False
    return r


def tb(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    return tf


def para(tf, runs, size=14, color=INK, align=PP_ALIGN.LEFT, bold=False,
         space_after=4, space_before=0, first=False, line=None):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    if space_after is not None:
        p.space_after = Pt(space_after)
    if space_before is not None:
        p.space_before = Pt(space_before)
    if line is not None:
        p.line_spacing = line
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, opts in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(opts.get("sz", size))
        r.font.bold = opts.get("bold", bold)
        r.font.color.rgb = opts.get("color", color)
        _set_font(r, latin=opts.get("lat", LAT_FONT), ea=opts.get("ea", EA_FONT))
    return p


def footer(slide, idx, total, dark=False):
    c = RGBColor(0x9A, 0xA3, 0xAF) if dark else SOFT
    tf = tb(slide, Inches(0.5), Inches(7.06), Inches(8.0), Inches(0.36))
    para(tf, [("模型／標準推導，全程非實測（modeled / standard-derived — never measured）",
               {"sz": 9, "color": c})], first=True, space_after=0)
    tf2 = tb(slide, Inches(11.6), Inches(7.06), Inches(1.2), Inches(0.36))
    para(tf2, [(f"{idx} / {total}", {"sz": 9, "color": c})],
         align=PP_ALIGN.RIGHT, first=True, space_after=0)


def title_bar(slide, zh, en):
    rect(slide, 0, 0, SW, Inches(1.06), PAPER)
    rect(slide, 0, Inches(1.06), SW, Pt(2.2), BRAND)
    rect(slide, Inches(0.5), Inches(0.30), Inches(0.12), Inches(0.52), BRAND)
    tf = tb(slide, Inches(0.78), Inches(0.20), Inches(11.9), Inches(0.78))
    para(tf, [(zh, {"sz": 25, "bold": True, "color": INK})], first=True, space_after=0)
    para(tf, [(en, {"sz": 12.5, "color": SOFT})], space_after=0)


TOTAL = 5


def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, DARKBG)
    rect(s, 0, Inches(3.30), SW, Pt(2.4), BRAND)
    tf = tb(s, Inches(0.9), Inches(2.05), Inches(11.5), Inches(1.25))
    para(tf, [("證據可讀性", {"sz": 44, "bold": True, "color": WHITE})], first=True, space_after=2)
    tf2 = tb(s, Inches(0.9), Inches(3.55), Inches(11.5), Inches(1.4))
    para(tf2, [("Evidence Legibility — 證明檢視器是即時運算，不是預錄動畫",
                {"sz": 18, "color": RGBColor(0xCF, 0xE6, 0xE2)})], first=True, space_after=6)
    para(tf2, [(f"示範航路 CHT × SANSA · {ds.WINDOW_LABEL_COMPACT}",
                {"sz": 13, "color": RGBColor(0x9A, 0xB6, 0xB2)})], space_after=2)
    tf3 = tb(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.1))
    para(tf3, [("觀眾證明：反應式即時重算　|　驗收委員證明：可稽核、可重現、誠實邊界",
                {"sz": 13, "color": RGBColor(0xCF, 0xE6, 0xE2)})], first=True, space_after=4)
    para(tf3, [("全程模型／標準推導，從不宣稱為實測或營運驗證；外部缺口僅揭露，不填造。",
                {"sz": 11, "color": RGBColor(0x8F, 0xA6, 0xA3)})], space_after=2)
    footer(s, 1, TOTAL, dark=True)


def tier_chip(slide, l, t, color, zh, en):
    rect(slide, l, t, Inches(0.16), Inches(0.30), color)
    tf = tb(slide, l + Inches(0.24), t - Inches(0.02), Inches(3.0), Inches(0.36), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(zh + "  ", {"sz": 11, "bold": True, "color": color}),
              (en, {"sz": 9, "color": SOFT})], first=True, space_after=0)


def simple_table(slide, l, t, col_w, headers, rows, row_h=Inches(0.34),
                 head_color=INK, head_bg=RGBColor(0xEE, 0xF1, 0xF4), zebra=True,
                 fontsz=11, cell_colors=None):
    x = l
    # header
    for i, hcell in enumerate(headers):
        rect(slide, x, t, col_w[i], row_h, head_bg)
        tf = tb(slide, x + Inches(0.06), t, col_w[i] - Inches(0.1), row_h, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, [(hcell, {"sz": fontsz, "bold": True, "color": head_color})], first=True, space_after=0)
        x += col_w[i]
    # rows
    y = t + row_h
    for ri, row in enumerate(rows):
        x = l
        rbg = RGBColor(0xFB, 0xFC, 0xFD) if (zebra and ri % 2) else PAPER
        for ci, cell in enumerate(row):
            rect(slide, x, y, col_w[ci], row_h, rbg, line=HAIR)
            col = INK
            if cell_colors and cell_colors.get((ri, ci)):
                col = cell_colors[(ri, ci)]
            tf = tb(slide, x + Inches(0.06), y, col_w[ci] - Inches(0.1), row_h, anchor=MSO_ANCHOR.MIDDLE)
            bold = ci == 0
            para(tf, [(str(cell), {"sz": fontsz, "color": col, "bold": bold,
                                   "lat": MONO if (ci > 0 and any(ch.isdigit() for ch in str(cell))) else LAT_FONT})],
                 first=True, space_after=0)
            x += col_w[ci]
        y += row_h


def slide_d1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, PAPER)
    title_bar(s, "D1 · 即時運算，不是預錄動畫", "Live compute, not a recording — move the rain slider, the numbers recompute")
    # left narrative
    tf = tb(s, Inches(0.6), Inches(1.5), Inches(5.1), Inches(4.8))
    para(tf, [("旋轉地球＋移動衛星，外觀與預錄動畫無異。", {"sz": 13, "color": INK})], first=True, space_after=8)
    para(tf, [("拖動雨量滑桿，吞吐量與抖動依 ", {"sz": 13, "color": INK}),
              ("ITU-R P.618-14", {"sz": 13, "bold": True, "color": STANDARD}),
              (" 即時重算 —", {"sz": 13, "color": INK})], space_after=2)
    para(tf, [("這是預錄動畫做不到的。", {"sz": 13, "bold": True, "color": BRAND})], space_after=10)
    para(tf, [("示範航路 GEO 代表幾何（el ≈ 31°，SGP4 傳播半徑）：",
               {"sz": 11.5, "color": SOFT})], space_after=4)
    para(tf, [("雨量 0 → 50 mm/h　吞吐 ", {"sz": 13, "color": INK}),
              ("43.9 → 11.2 Mbps", {"sz": 14, "bold": True, "color": PROXY, "lat": MONO}),
              ("（−74.4%）", {"sz": 12, "color": GAP})], space_after=2)
    para(tf, [("抖動 ", {"sz": 13, "color": INK}),
              ("8.0 → 18.9 ms", {"sz": 14, "bold": True, "color": PROXY, "lat": MONO})], space_after=12)
    para(tf, [("吞吐量＝模型容量代理（modeled capacity proxy），非封包實測；雨量為輸入值，非實測天氣。",
               {"sz": 10, "color": SOFT})], space_after=2)
    # right: sensitivity table
    tf2 = tb(s, Inches(6.05), Inches(1.42), Inches(6.6), Inches(0.4))
    para(tf2, [("雨量敏感度掃描（同一模型，由現場滑桿驅動）", {"sz": 12, "bold": True, "color": INK})],
         first=True, space_after=0)
    simple_table(
        s, Inches(6.05), Inches(1.85),
        [Inches(1.5), Inches(1.9), Inches(1.5), Inches(1.7)],
        ["雨量 Rain", "吞吐代理 Tput", "抖動 Jitter", "對比晴空"],
        [["0 mm/h", "43.9 Mbps", "8.0 ms", "reference"],
         ["5 mm/h", "33.8 Mbps", "10.1 ms", "−23.0%"],
         ["10 mm/h", "27.5 Mbps", "11.7 ms", "−37.2%"],
         ["25 mm/h", "18.2 Mbps", "15.0 ms", "−58.5%"],
         ["50 mm/h", "11.2 Mbps", "18.9 ms", "−74.4%"]],
        row_h=Inches(0.46), fontsz=12)
    tf3 = tb(s, Inches(6.05), Inches(4.7), Inches(6.6), Inches(1.6))
    para(tf3, [("來源 / Source", {"sz": 10, "bold": True, "color": SOFT})], first=True, space_after=2)
    para(tf3, [("報告「Audit & evidence」分頁 · 雨衰敏感度表（ITU-R P.618-14 §2.2.1.1）。"
                "現場示範對應 WS-A 重算提示與 docs/demo-reaction-script.md。",
                {"sz": 10, "color": SOFT})], space_after=2)
    footer(s, 2, TOTAL)


def pipe_box(slide, l, t, w, zh, std, tier_color, tier_label):
    h = Inches(1.55)
    rect(slide, l, t, w, h, PAPER, line=HAIR)
    rect(slide, l, t, w, Pt(3.2), tier_color)
    tf = tb(slide, l + Inches(0.08), t + Inches(0.12), w - Inches(0.16), h - Inches(0.2))
    para(tf, [(zh, {"sz": 12, "bold": True, "color": INK})], first=True, space_after=4, line=1.0)
    para(tf, [(std, {"sz": 8.5, "color": SOFT})], space_after=4, line=1.0)
    para(tf, [(tier_label, {"sz": 8.5, "bold": True, "color": tier_color})], space_after=0)


def slide_d2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, PAPER)
    title_bar(s, "D2 · 運算流程（每步標註標準條款）", "Compute pipeline — every step tagged with its standard")
    if os.path.exists(D2_FIG):
        # Vector pipeline rendered by research-visual-lab (editable SVG -> PNG).
        pic_w = Inches(11.6)
        s.shapes.add_picture(D2_FIG, Inches((13.333 - 11.6) / 2), Inches(1.35), width=pic_w)
    else:
        boxes = [
            ("公開 TLE\nCelesTrak GP", "NORAD · GP 公開資料", REAL, "公開來源"),
            ("SGP4 傳播\n→ ECI 位置/速度", "Vallado AIAA 2006 · WGS72", REAL, "即時運算"),
            ("ECI→ECEF→大地\nGMST 旋轉", "satellite.js · 互視幾何", REAL, "即時運算"),
            ("仰角 / 可視視窗\n互視交集", "ENU 仰角門檻", REAL, "即時運算"),
            ("鏈路預算\nFSPL·氣·雨·天線", "TR 38.811 §6.6.2 · P.676-13\nP.618-14 · S.1528/S.465", STANDARD, "標準模型"),
            ("換手決策\nargmax 代理", "TR 38.821 §7.3.2.2 + V-MO1", STANDARD, "標準模型"),
        ]
        n = len(boxes)
        gap_w = Inches(0.16)
        total_w = SW - Inches(1.0)
        box_w = (total_w - gap_w * (n - 1)) / n
        x = Inches(0.5)
        top = Inches(2.15)
        for i, (zh, std, tc, tl) in enumerate(boxes):
            pipe_box(s, x, top, box_w, zh, std, tc, tl)
            if i < n - 1:
                ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + box_w + Inches(0.005),
                                        top + Inches(0.62), gap_w - Inches(0.01), Inches(0.3))
                ar.fill.solid(); ar.fill.fore_color.rgb = HAIR; ar.line.fill.background(); ar.shadow.inherit = False
            x += box_w + gap_w
    # tiers legend + honesty note
    tier_chip(s, Inches(0.7), Inches(5.55), REAL, "即時運算", "Real compute (geometry/timing)")
    tier_chip(s, Inches(5.1), Inches(5.55), STANDARD, "標準模型", "Standard model (published formula)")
    tier_chip(s, Inches(9.7), Inches(5.55), PROXY, "示意代理", "Illustrative proxy (downstream)")
    tf = tb(s, Inches(0.7), Inches(6.05), Inches(12.0), Inches(0.95))
    para(tf, [("幾何／時序為即時重算的真實計算；鏈路量值為標準公式在代表幾何上的模型輸出；下游 吞吐／RSRP／抖動 為 ",
               {"sz": 11, "color": INK}),
              ("示意代理", {"sz": 11, "bold": True, "color": PROXY}),
              ("（無 EIRP／雜訊／頻寬、無封包實測）。", {"sz": 11, "color": INK})], first=True, space_after=5, line=1.1)
    para(tf, [("樣本（GEO，本航路）：FSPL d=38613.9 km / 20 GHz → 210.21 dB · 單向延遲 130.802 ms · 氣體吸收 0.539 dB · 合成天線增益 91.9 dB。",
               {"sz": 10, "color": SOFT})], space_after=2, line=1.1)
    footer(s, 3, TOTAL)


def slide_d3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, PAPER)
    title_bar(s, "D3 · 誠實分級矩陣", "Honesty grading matrix — what is real, modeled, proxy, or a disclosed gap")
    cards = [
        (REAL, "即時運算", "Real compute", "3 類",
         "衛星位置／軌跡 · 可視視窗／仰角 · 通訊時間／換手次數",
         "每航路即時重算；可由 Raw data 與來源溯源彈窗重現。"),
        (STANDARD, "標準模型", "Standard model", "5 類",
         "斜距／單向延遲 · FSPL · 雨衰 · 氣體吸收 · 天線增益",
         "TR 38.811 §6.6.2 · P.618-14 · P.676-13 · S.1528／S.465；公式真實、幾何為代表值，非實測。"),
        (PROXY, "示意代理", "Illustrative proxy", "3 類",
         "吞吐量 Mbps · RSRP · 抖動",
         "晴空容量基準經大氣／天線去化；無 EIRP／雜訊／頻寬，非封包實測速率。"),
        (GAP, "外部缺口", "External gap", "揭露待補",
         "封包實測 KPI · 原生 RF 換手 · 實站 RF 硬體 · 實測天氣／R0.01 · 實測 horizon · 驗收門檻",
         "檢視器範圍外；僅揭露，不造假填數。"),
    ]
    top = Inches(1.5)
    ch = Inches(1.30)
    for i, (color, zh, en, count, items, note) in enumerate(cards):
        y = top + i * (ch + Inches(0.07))
        rect(s, Inches(0.6), y, Inches(12.1), ch, PAPER, line=HAIR)
        rect(s, Inches(0.6), y, Inches(0.14), ch, color)
        tf = tb(s, Inches(0.95), y + Inches(0.08), Inches(2.7), ch - Inches(0.16), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, [(zh, {"sz": 16, "bold": True, "color": color})], first=True, space_after=1)
        para(tf, [(en, {"sz": 9.5, "color": SOFT})], space_after=1)
        para(tf, [(count, {"sz": 10.5, "bold": True, "color": INK})], space_after=0)
        tf2 = tb(s, Inches(3.75), y + Inches(0.10), Inches(8.85), ch - Inches(0.18), anchor=MSO_ANCHOR.MIDDLE)
        para(tf2, [(items, {"sz": 12, "bold": True, "color": INK})], first=True, space_after=3, line=1.0)
        para(tf2, [(note, {"sz": 10, "color": SOFT})], space_after=0, line=1.0)
    footer(s, 4, TOTAL)


def slide_d4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, PAPER)
    title_bar(s, "D4 · 驗證證據", "Verification evidence — re-runnable gates, retained checksums, independent cross-check")
    # left column: gates + reproducibility
    rect(s, Inches(0.6), Inches(1.5), Inches(5.95), Inches(5.25), PAPER, line=HAIR)
    tf = tb(s, Inches(0.85), Inches(1.65), Inches(5.5), Inches(5.0))
    para(tf, [("可重跑閘門 / Re-runnable gates", {"sz": 13, "bold": True, "color": BRAND})], first=True, space_after=6)
    for label, val in [
        ("Bucket-A 覆蓋閘 (g1)", "19 / 19 通過 · npm run verify:g1"),
        ("24 小時穩定 soak", "通過（保留紀錄 2026-05-16）"),
        ("過度宣稱 CI 閘", "通過 · npm test (verify-phase0)"),
        ("TLE 完整性／基線", "通過 · npm run verify:tle"),
        ("執行時物理單元測試", "通過 · npm run test:unit"),
    ]:
        para(tf, [("✔ ", {"sz": 12, "bold": True, "color": REAL}),
                  (label + "：", {"sz": 11.5, "color": INK}),
                  (val, {"sz": 10.5, "color": SOFT})], space_after=4, line=1.0)
    para(tf, [("可重現性 / Reproducibility", {"sz": 13, "bold": True, "color": BRAND})],
         space_after=6, space_before=8)
    para(tf, [("酬載校驗碼（FNV-1a）", {"sz": 11.5, "color": INK}),
              ("  f4e19f8b", {"sz": 12, "bold": True, "color": INK, "lat": MONO})], space_after=2, line=1.0)
    para(tf, [("相同航路＋時窗＋TLE 快照重跑即重現同一校驗碼。", {"sz": 10, "color": SOFT})], space_after=2, line=1.0)
    # right column: package + standards + sgp4
    rect(s, Inches(6.75), Inches(1.5), Inches(5.95), Inches(5.25), PAPER, line=HAIR)
    tf2 = tb(s, Inches(7.0), Inches(1.65), Inches(5.5), Inches(5.0))
    para(tf2, [("保留證據包 / Retained package", {"sz": 13, "bold": True, "color": BRAND})], first=True, space_after=6)
    para(tf2, [("報告 HTML　sha256 ", {"sz": 11, "color": INK}),
               ("9061837b…", {"sz": 11, "bold": True, "color": INK, "lat": MONO})], space_after=2, line=1.0)
    para(tf2, [("資料列 CSV　sha256 ", {"sz": 11, "color": INK}),
               ("e07359f8…", {"sz": 11, "bold": True, "color": INK, "lat": MONO})], space_after=2, line=1.0)
    para(tf2, [("Actor 63 列 · 可視 64 視窗 · 需求 34 列", {"sz": 10.5, "color": SOFT})], space_after=8, line=1.0)
    para(tf2, [("保留標準文件 / Retained standards", {"sz": 13, "bold": True, "color": BRAND})], space_after=6)
    para(tf2, [("12 份 3GPP / ITU-R PDF（sha256 列於套件 README）：", {"sz": 10.5, "color": SOFT})], space_after=2, line=1.0)
    para(tf2, [("TR 38.811 ", {"sz": 10, "color": INK}), ("824ea7…", {"sz": 10, "color": SOFT, "lat": MONO}),
               ("　P.618-14 ", {"sz": 10, "color": INK}), ("9812e7…", {"sz": 10, "color": SOFT, "lat": MONO})],
         space_after=1, line=1.0)
    para(tf2, [("P.676-13 ", {"sz": 10, "color": INK}), ("8c09b2…", {"sz": 10, "color": SOFT, "lat": MONO}),
               ("　S.1528 ", {"sz": 10, "color": INK}), ("a3b3d6…", {"sz": 10, "color": SOFT, "lat": MONO}),
               ("　S.465-6 ", {"sz": 10, "color": INK}), ("a813d8…", {"sz": 10, "color": SOFT, "lat": MONO})],
         space_after=8, line=1.0)
    para(tf2, [("SGP4 獨立交叉驗證 / Independent check", {"sz": 13, "bold": True, "color": BRAND})], space_after=6)
    para(tf2, [("獨立 python-sgp4（WGS72，Vallado AIAA 2006）對位：", {"sz": 10.5, "color": SOFT})], space_after=2, line=1.0)
    para(tf2, [("一致到 ", {"sz": 11, "color": INK}),
               ("0.057 m", {"sz": 13, "bold": True, "color": REAL, "lat": MONO}),
               ("，遠優於 1 km 容差 — 偽造者須湊出獨立 SGP4 也算得出的數字。", {"sz": 10.5, "color": INK})], space_after=2, line=1.0)
    footer(s, 5, TOTAL)


def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    slide_title(prs)
    slide_d1(prs)
    slide_d2(prs)
    slide_d3(prs)
    slide_d4(prs)
    prs.save(OUT)
    print("wrote", OUT, "(", len(prs.slides._sldIdLst), "slides )")


if __name__ == "__main__":
    main()
