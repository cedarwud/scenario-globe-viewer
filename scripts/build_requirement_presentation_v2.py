#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Requirement-coverage presentation on the repo's `template.pptx` (NTPU CSIE /
wMNLab academic theme), tuned for a procurement acceptance reviewer.

- Chinese categories 第一/二/三類 (no "Bucket" jargon).
- Chaptered: 14 thematic sub-sections (2-3 requirements each); status stated once
  per 本節狀態 header; every requirement row carries its own 邊界·不可宣稱.
- Tables at 16pt; any table that does not fit one page is auto-paginated (i/N).
- Uniform content top: every slide's body starts at the same height (just below
  the title rule), so no slide sits too low.
- Front control surface: coverage matrix, 不可宣稱/未證明事項 (plain list, no box),
  證據包與雜湊 traceability — before the per-requirement walk-through.
- Honesty: model / standards / geometric / assumed values are never presented as
  measured or operator-validated; external gaps disclosed-only; Tier A = 0.

Fonts: 繁中 = 標楷體 (DFKai-SB), English/code = Times New Roman. Title 28 / body 24.
Content reconciled with docs/requirement-presentation-walkthrough.md +
docs/data-source-index.md; canonical 34-row matrix / catalog imported from the v1
builder. v1 deck + original template.pptx untouched. Output:
`deliverable/requirement-presentation-2026-06-14-ntpu-template.pptx`.
"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import build_requirement_presentation as base   # noqa: E402
import viz_native   # noqa: E402  (native pptx renderer for box/flow *.viz.json)
import demo_scenario as ds   # noqa: E402  (demo window/route single source of truth)

ROOT = os.path.dirname(HERE)
TEMPLATE = os.path.join(ROOT, "template.pptx")
OUT = os.path.join(ROOT, "deliverable", "requirement-presentation-2026-06-14-ntpu-template.pptx")
FIGDIR = os.path.join(ROOT, "deliverable", "evidence-legibility-figures")
def figp(name): return os.path.join(FIGDIR, name)

EA = "標楷體"
LAT = "Times New Roman"

T_TITLE = 28
T_BODY = 24
T_SUB = 18
TBL_BODY = 16          # coverage-matrix cell text (kept 16: dense scan)
TBL_HEAD = 16
BIG_BODY = 20          # all other tables (group / gap / catalog / 證據包)
BIG_HEAD = 20
# uniform content top across all slides (just under the title rule)
TOP_STATUS = 1.20      # 本節狀態 line on group slides
TOP_GROUP_TABLE = 1.66 # group table (below 本節狀態)
TOP_CONTENT = 1.30     # text-box content slides
TOP_TABLE = 1.36       # table-only slides
PAGE_BOTTOM = 6.9      # keep tables/text above the template bottom rule

NAVY = RGBColor(0x2E, 0x2E, 0x6E)
INK = RGBColor(0x1A, 0x1A, 0x1A)
SOFT = RGBColor(0x5B, 0x63, 0x72)
DRED = RGBColor(0x7A, 0x2E, 0x2E)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
AMBER = RGBColor(0xE0, 0x8A, 0x1E)
REDST = RGBColor(0xC6, 0x28, 0x28)
GRAYST = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xF4, 0xF6, 0xF8)

CLS = base.CLS
ST = base.ST
ST_SHORT = {"complete": "完成", "partial": "部分(B)", "gap": "缺口", "scope": "範圍外"}
SHORT_CLS = {"public": "公開", "standards": "標準", "geometric": "幾何", "assumed": "假設",
             "display": "顯示", "external": "外部", "legacy": "快取", "gap": "缺口"}
CAT = {"A": ("第一類", "專案展示需求（demo 已實作）"),
       "B": ("第二類", "公開標準替代（待營運商 Tier A 資料）"),
       "C": ("第三類", "外部基礎設施 / 報告層（範圍外）")}

# ------------------------------------------------------------------ helpers
def setfont(run, latin=LAT, ea=EA):
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "zh-TW")
    rPr.set("altLang", "zh-TW")
    rPr.set("noProof", "1")
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", ea if tag == "a:ea" else latin)

def fit(tf):
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

def para(tf, runs, sz=T_BODY, align=PP_ALIGN.LEFT, sa=4, sb=0, line=1.1, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(sa); p.space_before = Pt(sb)
    if line:
        p.line_spacing = line
    pPr = p._p.get_or_add_pPr()
    for ex in pPr.findall(qn("a:buNone")): pPr.remove(ex)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))
    for t, o in runs:
        r = p.add_run(); r.text = t
        r.font.size = Pt(o.get("sz", sz)); r.font.bold = o.get("b", False)
        r.font.italic = o.get("i", False); r.font.color.rgb = o.get("c", INK)
        setfont(r, ea=o.get("ea", EA))
    return p

def fill_title(slide, runs):
    ph = slide.shapes.title
    tf = ph.text_frame; tf.clear(); fit(tf)
    p = tf.paragraphs[0]
    for t, o in runs:
        r = p.add_run(); r.text = t
        r.font.size = Pt(o["sz"]); r.font.bold = o.get("b", True)
        r.font.color.rgb = o.get("c", NAVY); setfont(r, ea=o.get("ea", EA))
    return ph

def clear_body_ph(slide):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == 0:
            continue
        if ph.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
            continue
        ph._element.getparent().remove(ph._element)

def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP, shrink=False):
    b = slide.shapes.add_textbox(l, t, w, h); tf = b.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    if shrink:
        fit(tf)
    return tf

def pic_caption(slide, path, l, t, w, caption=None):
    if not path or not os.path.exists(path):
        return None
    p = slide.shapes.add_picture(path, l, t, width=w)
    fr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, p.height)
    fr.fill.background(); fr.line.color.rgb = RGBColor(0xD9, 0xDD, 0xE3); fr.line.width = Pt(0.75)
    fr.shadow.inherit = False
    if caption:
        ct = textbox(slide, l, t + p.height + Inches(0.04), w, Inches(0.4))
        para(ct, [(caption, {"sz": 14, "c": SOFT, "i": True})], first=True, sa=0)
    return p

def style_cell(cell, text, sz=TBL_BODY, bold=False, color=INK, fill=WHITE, align=PP_ALIGN.LEFT):
    cell.margin_left = cell.margin_right = Inches(0.07)
    cell.margin_top = cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    p = cell.text_frame.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color
    setfont(r)

def _put_table(slide, headers, rows, col_w, top, bsz, hsz):
    nrows = len(rows) + 1
    h = min(PAGE_BOTTOM - top, 0.62 * nrows)
    gt = slide.shapes.add_table(nrows, len(headers), Inches(0.5), Inches(top), sum(col_w), Inches(h)).table
    gt.first_row = False; gt.horz_banding = False
    for j, w in enumerate(col_w):
        gt.columns[j].width = w
    for j, hd in enumerate(headers):
        style_cell(gt.cell(0, j), hd, sz=hsz, bold=True, color=WHITE, fill=NAVY)
    for i, row in enumerate(rows, 1):
        bandc = BAND if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            col, bold, fillc, al = INK, False, bandc, PP_ALIGN.LEFT
            if isinstance(val, tuple):
                val, meta = val
                col = meta.get("c", INK); bold = meta.get("b", False)
                fillc = meta.get("fill", bandc); al = meta.get("al", PP_ALIGN.LEFT)
            style_cell(gt.cell(i, j), val, sz=bsz, bold=bold, color=col, fill=fillc, align=al)
    return gt

def paginated_table(prs, L, title_main, sub, headers, rows, col_w, per_page, bsz=TBL_BODY, hsz=TBL_HEAD):
    nr = len(rows)
    pages = max(1, -(-nr // per_page))          # ceil; min 1 page
    base_sz, rem = (nr // pages, nr % pages) if nr else (0, 0)
    chunks, i = [], 0
    for p in range(pages):
        sz = base_sz + (1 if p < rem else 0)    # balanced sizes, no orphan page
        chunks.append(rows[i:i + sz]); i += sz
    n = len(chunks)
    for idx, chunk in enumerate(chunks, 1):
        s = prs.slides.add_slide(L[5])
        tr = [(title_main + (f"（{idx}/{n}）" if n > 1 else ""), {"sz": T_TITLE, "b": True, "c": NAVY})]
        if sub:
            tr.append(("　" + sub, {"sz": T_SUB, "b": True, "c": SOFT}))
        fill_title(s, tr); clear_body_ph(s)
        _put_table(s, headers, chunk, col_w, TOP_TABLE, bsz, hsz)

# ------------------------------------------------------------------ slide types
def s_title(prs, L):
    s = prs.slides.add_slide(L[0])
    fill_title(s, [("需求逐條對應 ｜ 資料來源誠實合約", {"sz": 32, "b": True, "c": NAVY})])
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame; tf.clear(); fit(tf)
            para(tf, [("LEO / MEO / GEO 衛星地面站場景檢視器 — 34 項需求", {"sz": 20, "c": NAVY})], first=True, sa=4)
            para(tf, [("「如何滿足 · 資料來源 · 引用 · 邊界 · 狀態」逐條對應", {"sz": 16, "c": SOFT})], sa=4)
            para(tf, [("資料來源分級：公開來源 · 標準模型推導 · 幾何推導 · 假設參數(Tier-B) · 顯示轉換 · 外部證據 · 舊快取 · 資料缺口",
                       {"sz": 14, "c": SOFT})], sa=4, line=1.15)
            para(tf, [("誠實鐵則：", {"sz": 15, "b": True, "c": DRED}),
                      ("模型 / 標準 / 幾何 / 假設 / 公開推導之值，絕不標示為「實測」或「營運商驗證」；外部缺口僅呈現為「已揭露待補」。",
                       {"sz": 14, "c": INK})], sa=4, line=1.2)
            para(tf, [("狀態日 2026-06-14　·　對齊 merge 1496a03（2026-06-13 公開資料缺口關閉）", {"sz": 12, "c": SOFT})], sa=0)
    return s

def s_agenda(prs, L):
    s = prs.slides.add_slide(L[1])
    fill_title(s, [("大綱", {"sz": T_TITLE, "b": True, "c": NAVY}), ("　Agenda", {"sz": T_SUB, "b": True, "c": SOFT})])
    clear_body_ph(s)
    lc = textbox(s, Inches(0.6), Inches(TOP_CONTENT), Inches(6.1), Inches(5.2), shrink=True)
    rc = textbox(s, Inches(6.9), Inches(TOP_CONTENT), Inches(5.9), Inches(5.2), shrink=True)
    left_items = ["1　誠實閱讀規則與來源分級", "2　需求總覽（第一/二/三類 = 19/7/8）", "3　選定配對 demo route",
                  "4　驗收總覽 · 34 列覆蓋矩陣", "5　不可宣稱 / 未證明事項"]
    right_items = ["6　證據包（可證 / 不可證）", "7　第一類：專案展示需求（8 小節）", "8　第二類：公開標準替代（3 小節）",
                   "9　第三類：外部 / 報告層（3 小節）", "10　資料缺口總表 · 來源目錄 · 誠實小結"]
    for i, it in enumerate(left_items):
        para(lc, [(it, {"sz": T_BODY, "b": True, "c": NAVY})], first=(i == 0), sa=13, line=1.2)
    for i, it in enumerate(right_items):
        para(rc, [(it, {"sz": T_BODY, "b": True, "c": NAVY})], first=(i == 0), sa=13, line=1.2)
    return s

def s_rule(prs, L):
    s = prs.slides.add_slide(L[1])
    fill_title(s, [("誠實閱讀規則", {"sz": T_TITLE, "b": True, "c": NAVY}), ("　每個數字問兩個問題", {"sz": T_SUB, "b": True, "c": SOFT})])
    clear_body_ph(s)
    lt = textbox(s, Inches(0.55), Inches(TOP_CONTENT), Inches(6.55), Inches(5.4), shrink=True)
    para(lt, [("1　顯示值從何而來？", {"sz": T_BODY, "b": True, "c": NAVY})], first=True, sa=2)
    para(lt, [("（公開產物 / 標準模型 / TLE 幾何 / 缺口）", {"sz": 14, "c": SOFT})], sa=9)
    para(lt, [("2　該來源實際能支撐什麼主張？", {"sz": T_BODY, "b": True, "c": NAVY})], sa=11)
    para(lt, [("功能展示 ≠ 來源權威。", {"sz": 22, "b": True, "c": DRED})], sa=6, line=1.18)
    para(lt, [("模型 / 假設 / 公開推導值不呈現為實測或營運商驗證；外部缺口僅「已揭露待補」。Tier A 營運商驗證 = 0 交付。",
               {"sz": 15, "c": INK})], sa=11, line=1.22)
    para(lt, [("標示說明：", {"sz": 15, "b": True, "c": NAVY})], sa=4)
    para(lt, [("本節狀態 = 驗收狀態　·　第一/二/三類 = 合約需求分組　·　來源類別 = 證據權威", {"sz": 14, "c": SOFT})], sa=0, line=1.2)
    rt = textbox(s, Inches(7.25), Inches(TOP_CONTENT), Inches(5.6), Inches(5.4), shrink=True)
    para(rt, [("資料來源分類（8 類 · 證據權威）", {"sz": T_SUB, "b": True, "c": NAVY})], first=True, sa=9)
    for k in ["public", "standards", "geometric", "assumed", "display", "external", "legacy", "gap"]:
        zh, en, col = CLS[k]
        para(rt, [("■ ", {"sz": 18, "c": col}), (f"{zh}　", {"sz": 17, "b": True, "c": col}), (en, {"sz": 13, "b": True, "c": col})], sa=9)
    return s

def s_overview(prs, L):
    s = prs.slides.add_slide(L[1])
    fill_title(s, [("需求總覽", {"sz": T_TITLE, "b": True, "c": NAVY}), ("　34 列 · 第一/二/三類 = 19 / 7 / 8", {"sz": T_SUB, "b": True, "c": SOFT})])
    clear_body_ph(s)
    lt = textbox(s, Inches(0.55), Inches(TOP_CONTENT), Inches(6.6), Inches(5.3), shrink=True)
    para(lt, [("第一類 · 專案展示需求（19）", {"sz": T_BODY, "b": True, "c": NAVY})], first=True, sa=3)
    para(lt, [("Cesium 球體 + 選定配對 V4 投影；幾何 / 標準推導，顯示轉換。「完成」= 模型 / 展示邊界內完成，非合約來源權威。", {"sz": 15, "c": INK})], sa=11, line=1.22)
    para(lt, [("第二類 · 公開標準替代（7）", {"sz": T_BODY, "b": True, "c": NAVY})], sa=3)
    para(lt, [("各列需營運商資料達 Tier A；以 ITU-R / 3GPP 公開標準自實作替代（4 已 runtime-wired，3 irreducible）。", {"sz": 15, "c": INK})], sa=11, line=1.22)
    para(lt, [("第三類 · 外部基礎設施 / 報告層（8）", {"sz": T_BODY, "b": True, "c": NAVY})], sa=3)
    para(lt, [("不在選定配對 viewer 範圍：6 外部網路 / 測試基礎設施 + 2 報告層交付。", {"sz": 15, "c": INK})], sa=15, line=1.22)
    para(lt, [("34", {"sz": 34, "b": True, "c": NAVY}), ("　需求　　", {"sz": 16, "c": INK}),
              ("0", {"sz": 34, "b": True, "c": DRED}), ("　Tier A 交付", {"sz": 16, "c": INK})], sa=0)
    pic_caption(s, base.ch("chart_buckets.png"), Inches(7.35), Inches(1.55), Inches(5.5), "需求分布（共 34 列；第一/二/三類）")
    return s

def s_demo(prs, L):
    s = prs.slides.add_slide(L[1])
    fill_title(s, [("選定配對 demo route", {"sz": T_TITLE, "b": True, "c": NAVY}), ("　CHT Yangmingshan ↔ SANSA", {"sz": T_SUB, "b": True, "c": SOFT})])
    clear_body_ph(s)
    lt = textbox(s, Inches(0.55), Inches(TOP_CONTENT), Inches(5.35), Inches(5.3), shrink=True)
    para(lt, [("路由：", {"sz": 16, "b": True, "c": NAVY})], first=True, sa=3)
    para(lt, [(ds.DEMO_ROUTE,
               {"sz": 14, "c": NAVY})], sa=10, line=1.2)
    para(lt, [("視窗：", {"sz": 16, "b": True, "c": NAVY}), (ds.WINDOW_LABEL_ZH, {"sz": 15, "c": INK})], sa=8, line=1.18)
    para(lt, [("預錄 TLE：", {"sz": 16, "b": True, "c": NAVY}), ("OneWeb/Galileo/GEO 2026-05-13~15 釘死快照，視窗在其 epoch 後 +1.8~5.7 天（近-epoch SGP4）。計分板「2026-06-13」屬另一條目錄新鮮度/規模證明，非此 demo 幾何來源。", {"sz": 13, "c": INK})], sa=8, line=1.2)
    para(lt, [("來源 tier：", {"sz": 16, "b": True, "c": NAVY}), ("geometric-derived；cross-family-geometric。", {"sz": 15, "c": INK})], sa=8, line=1.18)
    para(lt, [("標籤掃描：", {"sz": 16, "b": True, "c": NAVY}), ("HTML 報告 39 standards-derived + 12 geometric-derived chip，0 operator-validated / public-disclosed。", {"sz": 15, "c": INK})], sa=8, line=1.2)
    para(lt, [("即：產生器不將任何模型 / 幾何值升級為營運商權威。", {"sz": 15, "b": True, "c": GREEN})], sa=0, line=1.18)
    pic_caption(s, base.IMG_GLOBE, Inches(6.1), Inches(TOP_CONTENT + 0.1), Inches(6.8), "viewer UI 範例 · 選定配對端點 · 模型換手（retained screenshot）")
    return s

def s_data_model(prs, L):
    s = prs.slides.add_slide(L[1])
    fill_title(s, [("資料更新機制", {"sz": T_TITLE, "b": True, "c": NAVY}), ("　兩條 TLE · 維運指令", {"sz": T_SUB, "b": True, "c": SOFT})])
    clear_body_ph(s)
    tb = textbox(s, Inches(0.5), Inches(TOP_CONTENT), Inches(12.33), Inches(2.05), shrink=True)
    para(tb, [("本檢視器有兩條獨立的 TLE 資料線，回答不同問題：", {"sz": 16, "b": True, "c": NAVY})], first=True, sa=5)
    para(tb, [("① Demo 幾何（預錄 · 可重現）", {"sz": 15, "b": True, "c": INK}), ("　釘死快照、視窗近其 epoch、固定不動（一旦變動就無法重現、與保存證據對不上）＝ 驗收場景基準日。", {"sz": 14, "c": INK})], sa=4, line=1.18)
    para(tb, [("② Catalog（新鮮度證明）", {"sz": 15, "b": True, "c": INK}), ("　定期更新到最新 CelesTrak，證明系統能即時吃進當前資料、≥500 LEO。", {"sz": 14, "c": INK})], sa=4, line=1.18)
    para(tb, [("　固定快照日 ＝ 驗收常態且誠實；簡報為未追蹤產物，隨時可重建（讀 config，日期自動正確）。", {"sz": 13, "b": True, "c": GREEN})], sa=0, line=1.15)
    headers = ["指令", "做什麼", "動到的資料", "何時用"]
    rows = [
        [("npm run update:tle", {"c": NAVY, "b": True}), "更新 catalog 到最新（CelesTrak）；不碰 demo", "catalog 新鮮度（leo-latest…）", "證明能吃當前資料"],
        [("npm run update:demo", {"c": NAVY, "b": True}), "照 config 重生 demo 衍生產物（golden · 證據）＋ 驗證；不換資料", "demo 保存證據 · golden", "改模型 / 設定後重對齊"],
        [("npm run repin:demo", {"c": NAVY, "b": True}), "換 demo 釘死快照為 catalog 最新 ＋ 移視窗 ＋ update:demo", "demo 幾何（換新資料）", "要 demo 改用最新資料"],
        [("npm run update:all", {"c": NAVY, "b": True}), "update:tle ＋ update:demo", "兩者", "例行同步"],
    ]
    col_w = [Inches(2.95), Inches(4.35), Inches(2.55), Inches(2.45)]
    _put_table(s, headers, rows, col_w, 3.55, 13, TBL_HEAD)
    return s

def _ops_slide(prs, L, title_main, sub, cmd, does, warn, meaning):
    s = prs.slides.add_slide(L[1])
    fill_title(s, [(title_main, {"sz": T_TITLE, "b": True, "c": NAVY}), ("　" + sub, {"sz": T_SUB, "b": True, "c": SOFT})])
    clear_body_ph(s)
    lt = textbox(s, Inches(0.6), Inches(TOP_CONTENT + 0.25), Inches(12.1), Inches(4.9), shrink=True)
    para(lt, [("指令　", {"sz": 17, "b": True, "c": NAVY}), (cmd, {"sz": 17, "b": True, "c": NAVY})], first=True, sa=14)
    para(lt, [("做什麼　", {"sz": 16, "b": True, "c": INK}), (does, {"sz": 16, "c": INK})], sa=14, line=1.28)
    para(lt, [("⚠ 注意　", {"sz": 16, "b": True, "c": DRED}), (warn, {"sz": 16, "c": INK})], sa=14, line=1.3)
    para(lt, [("數據意義　", {"sz": 16, "b": True, "c": GREEN}), (meaning, {"sz": 16, "c": INK})], sa=0, line=1.28)
    return s

def s_ops_tle(prs, L):
    return _ops_slide(prs, L,
        "維運指令 ①", "npm run update:tle · 更新 catalog 新鮮度",
        "npm run update:tle",
        "抓 CelesTrak 最新 GP/TLE 到 catalog（satellites-network/）。這是「系統能即時吃進當前衛星資料」的能力證明。",
        "CelesTrak 官方限流：同一筆資料約 2 小時內勿重複抓取，否則會被擋（回 error）。例行更新即足夠；refresh:tle 內建 --if-older-than-days 7 門檻，自動避免過度抓取。",
        "catalog = 新鮮度／規模證明（≥500 LEO），不是 demo 幾何來源（demo 用釘死快照）。")

def s_ops_cron(prs, L):
    s = prs.slides.add_slide(L[1])
    fill_title(s, [("維運指令 ②", {"sz": T_TITLE, "b": True, "c": NAVY}), ("　TLE 排程自動更新（cron）", {"sz": T_SUB, "b": True, "c": SOFT})])
    clear_body_ph(s)
    tb = textbox(s, Inches(0.6), Inches(TOP_CONTENT), Inches(12.1), Inches(1.85), shrink=True)
    para(tb, [("做什麼　", {"sz": 16, "b": True, "c": INK}), ("把 catalog 更新排成 cron，每天自動跑兩次（08:00 / 20:00 伺服器本地時間），免手動。每輪＝refresh:tle ＋ refresh:tle:local（匯入本地快取）。", {"sz": 15, "c": INK})], first=True, sa=8, line=1.25)
    para(tb, [("⚠ 注意　", {"sz": 15, "b": True, "c": DRED}), ("有 lock 防併發；log（logs/tle/）超過 1MB 自動壓縮、保留 30 天；仍受 CelesTrak ~2h 限流（refresh:tle 門檻自動保護）。demo 用釘死快照，不受排程影響。", {"sz": 14, "c": INK})], sa=0, line=1.25)
    headers = ["指令", "做什麼"]
    rows = [
        [("npm run tle:schedule:install", {"c": NAVY, "b": True}), "安裝 cron 排程（每天 08:00、20:00 自動更新 catalog）"],
        [("npm run tle:schedule:status", {"c": NAVY, "b": True}), "查排程是否安裝 ＋ log 狀態"],
        [("npm run tle:schedule:logs", {"c": NAVY, "b": True}), "看最近更新紀錄"],
        [("npm run tle:schedule:run", {"c": NAVY, "b": True}), "立刻手動跑一輪（不等排程）"],
        [("npm run tle:schedule:uninstall", {"c": NAVY, "b": True}), "移除排程"],
    ]
    col_w = [Inches(4.7), Inches(7.6)]
    _put_table(s, headers, rows, col_w, 3.25, 14, TBL_HEAD)
    return s

def s_ops_gate(prs, L):
    return _ops_slide(prs, L,
        "維運指令 ③", "npm run verify:tle · 驗證 gate 執行環境",
        "npm run verify:tle（＋ G1 / G3 / G4）",
        "在無頭瀏覽器驗證 runtime 投影、可見視窗、換手等數字是否符合保存基準。",
        "這些 gate 在瀏覽器直接 import 原始碼，需在「乾淨環境 / CI」啟動 dev-server 執行。本機長時 dev-server 對 session 內新建的 public 檔有服務怪癖（curl／build 版正常，僅瀏覽器 fetch 偶 404）—— re-pin 後於乾淨 shell 重跑即可，非資料問題。",
        "gate 驗的是 demo 幾何數字；換資料（repin:demo）後需重新 baseline（人工確認新數字合理）。")

def s_donotclaim(prs, L):
    s = prs.slides.add_slide(L[1])
    fill_title(s, [("不可宣稱 / 未證明事項", {"sz": T_TITLE, "b": True, "c": NAVY}), ("　已揭露待補（外部 / 營運商資料）", {"sz": T_SUB, "b": True, "c": DRED})])
    clear_body_ph(s)
    tf = textbox(s, Inches(0.6), Inches(TOP_CONTENT), Inches(12.2), Inches(5.4), shrink=True)
    para(tf, [("本簡報絕不宣稱以下任一項 — 皆為已揭露之資料缺口：", {"sz": 18, "b": True, "c": DRED})], first=True, sa=10, line=1.2)
    items = [
        "封包延遲 / 抖動 / 吞吐 / 遺失（無 ping/iperf 封包測試）",
        "原生 RF 換手 / 主動服務路徑 / 主動服務衛星 / gateway 指派",
        "營運商站台 RF 硬體真值（EIRP、天線盤徑、極化、波束、天線型號）",
        "勘測 RF 視界 / 全 69 列 DEM 替換",
        "實測連線天氣 / 長期 R0.01 可用率校準",
        "外部驗收門檻腳本 / 最終書面報告 / WP1 評估報告",
        "實體網路橋接 / NAT / tunnel / DUT / 流量產生器驗證",
    ]
    for it in items:
        para(tf, [("✕　", {"sz": 17, "b": True, "c": DRED}), (it, {"sz": 17, "c": INK})], sa=8, line=1.18)
    para(tf, [("Tier A 營運商驗證 = 0 交付；全部第二類列待 Tier A 換入。功能展示 ≠ 來源權威。",
               {"sz": 16, "b": True, "c": NAVY})], sb=4, sa=0, line=1.2)
    return s

def s_section(prs, L, bucket):
    cn, desc = CAT[bucket]
    s = prs.slides.add_slide(L[1])
    fill_title(s, [(f"{cn} · {desc}", {"sz": T_TITLE, "b": True, "c": NAVY})])
    clear_body_ph(s)
    tf = textbox(s, Inches(0.6), Inches(TOP_CONTENT), Inches(11.9), Inches(5.0), shrink=True)
    body = {
        "A": ["主軸：Cesium 球體 viewer + 選定配對 V4 地面站投影。",
              "誠實來源類別：幾何推導（TLE + SGP4 + 可見性）+ 標準模型推導（link budget），經顯示轉換（面板、報告、CSV）。",
              "本類「完成」= 功能 / 模型 / 展示邊界內完成；無任何列宣稱主動服務路徑、封包量測或營運商遙測。"],
        "B": ["各列需營運商資料才能達 Tier A 權威；本專案於 src/runtime/link-budget/ 自實作 Tier B 公開標準替代，附保存之 ITU-R / 3GPP 標準 PDF。",
              "4 替代已接入選定配對 runtime；3 條 irreducible 有合理公開替代但其最強用語仍為缺口。",
              "2026-06-13 merge 改變三列：天線已 runtime-wired、已有公開在地雨統計、選定配對已採 DEM / 地形 provenance。"],
        "C": ["不在選定配對 viewer 之 UI 範圍。6 列為本專案以外擁有之外部網路 / 測試基礎設施；2 列為本專案另行撰寫之報告層交付。",
              "列出以求需求清單完整。選定配對投影為球體 / runtime 模型產物，無法自證以下任一列。"],
    }[bucket]
    for i, ln in enumerate(body):
        para(tf, [(ln, {"sz": T_BODY, "c": INK})], first=(i == 0), sa=11, line=1.28)
    return s

def s_group(prs, L, g):
    cn = CAT[g["bucket"]][0]
    s = prs.slides.add_slide(L[5])
    fill_title(s, [(f"{cn} ▸ {g['title']}", {"sz": T_TITLE, "b": True, "c": NAVY})])
    clear_body_ph(s)
    tag = textbox(s, Inches(0.55), Inches(TOP_STATUS), Inches(12.2), Inches(0.4))
    para(tag, [("本節狀態：", {"sz": 16, "b": True, "c": SOFT}), (g["status_text"], {"sz": 16, "b": True, "c": g["status_color"]})], first=True, sa=0)
    headers = ["需求", "滿足 / 模組", "來源類別", "邊界 · 不可宣稱", "需求 ID"]
    col_w = [Inches(2.8), Inches(2.9), Inches(1.3), Inches(3.35), Inches(1.8)]
    rows = []
    for (rid, need, sat, clskeys, exc, bnd) in g["rows"]:
        need_cell = (need, {})
        if exc:
            ecol, _ = ST[exc]; emark = {"partial": "（部分）", "gap": "（缺口）", "scope": "（範圍外）"}[exc]
            need_cell = (need + emark, {"c": ecol, "b": True})
        rows.append([
            need_cell,
            sat,
            ("／".join(SHORT_CLS[k] for k in clskeys), {"b": True, "c": CLS[clskeys[0]][2]}),
            (bnd, {"c": DRED}),
            (rid, {"c": SOFT}),
        ])
    _put_table(s, headers, rows, col_w, TOP_GROUP_TABLE, BIG_BODY, BIG_HEAD)
    return s

def s_modelcharts(prs, L):
    s = prs.slides.add_slide(L[1])
    fill_title(s, [("第二類 · 模型輸出（範例）", {"sz": T_TITLE, "b": True, "c": NAVY}), ("　標準 + 假設推導 · 非實測", {"sz": T_SUB, "b": True, "c": SOFT})])
    clear_body_ph(s)
    pic_caption(s, base.ch("chart_antenna.png"), Inches(0.7), Inches(TOP_CONTENT + 0.15), Inches(5.9), "合併天線增益 S.1528 / S.465 + 假設參數 · RSRP 相對 proxy")
    pic_caption(s, base.ch("chart_rain.png"), Inches(6.85), Inches(TOP_CONTENT + 0.15), Inches(5.9), "降雨衰減 P.618-14 · 範例 what-if 輸入 · 非實測")
    bn = textbox(s, Inches(0.6), Inches(6.2), Inches(12.0), Inches(0.75), shrink=True)
    para(bn, [("邊界　", {"sz": 14, "b": True, "c": DRED}),
              ("標準模型推導 + 假設天線參數；雨為互動 what-if 輸入，晴空報告 rain=0。非實測天氣、非營運商 RF 硬體真值。", {"sz": 14, "c": DRED})], first=True, sa=0, line=1.14)
    return s

def s_closing(prs, L):
    s = prs.slides.add_slide(L[1])
    fill_title(s, [("誠實小結", {"sz": T_TITLE, "b": True, "c": NAVY})])
    clear_body_ph(s)
    tf = textbox(s, Inches(0.6), Inches(TOP_CONTENT), Inches(12.0), Inches(5.2), shrink=True)
    rows = [
        ("功能展示 ≠ 來源權威。", "多數 UI / 模型需求已在 demo route 功能展示；封包 / RF / 天氣 / 硬體證據仍受邊界限制或缺口。", NAVY),
        ("Tier A 營運商驗證 = 0 交付。", "全部第二類列均待 Tier A 換入；目前以 src/runtime/link-budget/ 之 Tier B 公開標準替代。", NAVY),
        ("本簡報絕不宣稱：", "封包延遲 / 抖動 / 吞吐 / 遺失 · 原生 RF 換手 / 主動服務路徑 · 營運商站台 RF 硬體真值 · 實測當地天氣 / R0.01 校準 · 外部驗收門檻 · 最終書面報告。", DRED),
        ("可宣稱：", "公開來源可追溯 · 標準模型依揭露公式 · 幾何 / 可見性 / 模型選出換手候選 · 假設參數明確標示 · 24h soak 保存產物。", GREEN),
    ]
    for i, (h, b, c) in enumerate(rows):
        para(tf, [(h + " ", {"sz": T_BODY, "b": True, "c": c}), (b, {"sz": 18, "c": INK})], first=(i == 0), sa=13, line=1.22)
    para(tf, [("逐列來源 / 引用 / 邊界見 docs/requirement-presentation-walkthrough.md　·　缺口登記見 docs/data-source-index.md",
               {"sz": 13, "c": SOFT})], sa=0, line=1.15)
    return s

def strip_hash(t):
    """Remove checksum / hash mentions (user: nobody reads SHA values) + leftover jargon."""
    for s in ["（sha d184c8dc…）", "（sha 662b9238…）", " + sha ×2", "+ sha ×2",
              "、檔案雜湊", "、檔案校驗碼", "、checksum", " + checksum", "+ checksum"]:
        t = t.replace(s, "")
    return t.replace("雜湊", "校驗碼").replace("bucket", "分類")

def cov_rows(reqs):
    out = []
    for r in reqs:
        col, _ = ST[r["status"]]
        out.append([r["name"], (ST_SHORT[r["status"]], {"c": col, "b": True}), base.cls_short(r["classes"]), (r["id"], {"c": SOFT})])
    return out

TRACE_HEADERS = ["證據", "可證", "不可證"]
TRACE_ROWS = [
    ["報告 HTML（runtime-projection-evidence-…-360m.html）", "需求 / 來源 / 模型 / 換手 / 可見性可讀報告", ("原始證據庫、重跑、封包/RF、驗收", {"c": DRED})],
    ["列級 CSV（runtime-projection-…-360m.csv）", "列級欄位 / 來源 tier / 模型輸出 / 非主張", ("payload 以外敘事、封包、驗收", {"c": DRED})],
    ["24h soak（summary.json · passed · 0 fail）", "保存 artifact 之 24 小時穩定", ("驗收環境穩定、重跑", {"c": DRED})],
    ["保存包（selected-pair-source-evidence/）", "選定配對來源包、route、產生時間", ("封包/RF/營運商驗證、最終驗收", {"c": DRED})],
    ["標準 PDF（3gpp-itu-references/）", "ITU-R / 3GPP §-引用對照、標準模型依據", ("營運商資料、實測品質、硬體", {"c": DRED})],
    ["公開資料（CWA rain · Copernicus DEM · registry）", "公開來源可追溯（雨 / 高程 / 站台）", ("實測連線天氣、R0.01、勘測視界", {"c": DRED})],
]

# ------------------------------------------------------------------ thematic sub-sections (chaptered; per-row boundary)
# rows: (ID, 需求, 滿足/模組, [source-class keys], exception status|None, 邊界·不可宣稱)
GROUPS = [
 dict(bucket="A", title="軌道模型整合與多軌切換", status_text="完成 · 邊界內", status_color=GREEN, rows=[
   ("R1-T1 / K-A1", "整合三軌 + 多軌切換", "SGP4(manifest)→handover-policy.ts", ["geometric"], None, "模型切換候選；非主動路徑 / 營運商記錄"),
   ("R1-T2 / K-D1", "整合衛星軌道資料", "CelesTrak GP/TLE manifest", ["public"], None, "快照；非即時目錄 / 營運商指派"),
   ("R1-F1 / K-E1", "≥500 LEO actor", "651 OneWeb LEO（全 1258 含 MEO/GEO）", ["public", "geometric"], None, "≥500 由 OneWeb；非 Starlink 預設"),
 ]),
 dict(bucket="A", title="TLE 追蹤與軌道里程碑", status_text="完成 · 邊界內", status_color=GREEN, rows=[
   ("K-A4", "輸入 TLE + 追蹤衛星", "resolveTleDate；晶片 2026-05-31", ["public"], None, "讀打包檔；非即時目錄連線"),
   ("R1-D1", "11/30 里程碑：軌道匯入", "匯入 artifacts + 傳播 actor", ["public", "geometric"], None, "不自證完整里程碑包"),
 ]),
 dict(bucket="A", title="換手與跨軌", status_text="完成 · 邊界內", status_color=GREEN, rows=[
   ("R1-T5 / K-D4", "換手規則 + 可調參數", "handover-policy.ts(TR 38.821 §7.3.2.2)", ["standards"], None, "模型控制；非 SLA / 封包軌跡"),
   ("R1-F4 / K-E4 / K-F4", "依延遲/抖動/速率切換", "link-budget → 換手引擎", ["standards"], None, "無營運商事件 / RF 換手軌跡"),
   ("V-MO1", "跨軌 LIVE 換手", "demo-balanced-v1（demo route 採用）", ["geometric"], None, "模型策略遷移;非原生 RF 換手;此 pair LEO=0，實為 GEO⇄MEO"),
 ]),
 dict(bucket="A", title="連結預算與通訊速率", status_text="完成 · 邊界內", status_color=GREEN, rows=[
   ("R1-T6 / K-D5", "通訊速率視覺化", "link-budget → 吞吐估計", ["standards", "assumed"], None, "延遲/吞吐為模型估計；RF 硬體缺口"),
   ("K-E6", "降雨衰減影響", "rain-attenuation.ts(P.618-14)", ["standards"], None, "what-if 雨率；非實測天氣"),
 ]),
 dict(bucket="A", title="可通訊時間與統計", status_text="多為完成；R1-F3 部分", status_color=AMBER, rows=[
   ("R1-F3 / K-E3", "即時可通訊時間", "可見視窗 + 連線選擇策略", ["geometric", "standards"], "partial", "模型可用時間；非 ping/iperf 實測"),
   ("R1-D3", "11/30 里程碑：通訊統計", "communicationStats payload", ["geometric"], None, "模型推導；非實測吞吐"),
 ]),
 dict(bucket="A", title="視覺場景與互動 UI", status_text="完成 · 邊界內", status_color=GREEN, rows=[
   ("R1-T3 / K-D2", "等效視覺場景(Cesium)", "Cesium 球體 + 報告", ["display"], None, "顯示轉換；不生新來源真實性"),
   ("R1-T4 / K-D3", "互動式 UI", "V4 側欄 + 報告分頁 + CSV", ["display"], None, "非完整 UI 互動軌跡"),
 ]),
 dict(bucket="A", title="報告匯出與 demo", status_text="完成 · 邊界內", status_color=GREEN, rows=[
   ("R1-F5 / K-E5", "統計報告匯出", "HTML 報告 + 列級 CSV 匯出", ["display", "external"], None, "證據面；非最終驗收證明"),
   ("K-F7", "產出 demo", "選定配對 route + 證據包", ["display", "external"], None, "代表性 demo；非外部驗收門檻"),
 ]),
 dict(bucket="A", title="模擬控制與穩定性里程碑", status_text="完成 · 邊界內", status_color=GREEN, rows=[
   ("R1-F2 / K-E2", "可調速度 30/60/120×", "m8a-v4-product-ux-model.ts；360min", ["display"], None, "app 介面；非列級 CSV 指標"),
   ("R1-D2", "11/30 里程碑：動態參數 UI", "側欄保留視窗/雨/策略", ["display"], None, "快照；非完整控制記錄"),
   ("R1-D4", "11/30 里程碑：24h 穩定", "soak passed · 86400000ms · 0 fail", ["external"], None, "保存產物；開報告不重跑 soak"),
 ]),
 dict(bucket="B", title="連線品質與天線（Tier B）", status_text="部分 · Tier B（待 Tier A）", status_color=AMBER, rows=[
   ("K-A2", "連線品質規則(延遲/抖動)", "TR 38.811 §6.6.2 + 5.3.1.1", ["standards"], None, "模型錨點；非 ping/RTT/jitter 量測"),
   ("K-A3-a", "天線參數(增益/波束/場型)", "S.1528 + S.465；接入 ln610/616", ["standards", "assumed"], None, "RSRP 相對 proxy；非營運商 RF 硬體"),
 ]),
 dict(bucket="B", title="降雨與整合（Tier B）", status_text="部分 · Tier B（待 Tier A）", status_color=AMBER, rows=[
   ("K-A3-b", "降雨衰減模型", "P.618-14 §2.2.1.1 全 path", ["standards"], None, "瞬時雨率；非 R0.01 / 實測天氣"),
   ("K-F2", "整合 V-team 模擬(替代)", "五個 link-budget 模組全接入", ["standards", "assumed"], None, "替代非 V-team 本身；缺站台硬體"),
 ]),
 dict(bucket="B", title="irreducible 缺口", status_text="irr-3 部分；irr-1/-2 缺口", status_color=REDST, rows=[
   ("irreducible-1", "真實封包軌跡", "僅 TR 38.811 模型基線", ["standards"], "gap", "無封包軌跡；不得宣稱實測延遲/抖動"),
   ("irreducible-2", "外部驗收腳本 + 門檻", "五條代表性 baseline route", ["external"], "gap", "代表 route；非驗收門檻腳本"),
   ("irreducible-3", "台灣在地雨統計校準", "CWA 46691/46693 → 5 mm/h", ["public", "assumed"], "partial", "公開樣本；非 measured-for-link / R0.01"),
 ]),
 dict(bucket="C", title="外部網路橋接與環境", status_text="範圍外（外部基礎設施）", status_color=GRAYST, rows=[
   ("K-F1", "實體網路橋接(INET↔ESTNeT)", "外部基礎設施", ["gap"], None, "模型可見性無法證明實體橋接"),
   ("K-A5", "Linux + Windows + WSL 環境", "外部驗證", ["external"], None, "非 UI / runtime 投影值"),
   ("K-B1", "Windows tunneling 真流量", "拓樸圖", ["gap"], None, "無 tunnel / 真流量橋接證據"),
 ]),
 dict(bucket="C", title="路由 / DUT / 流量產生器", status_text="範圍外（外部基礎設施）", status_color=GRAYST, rows=[
   ("K-C1", "INET NAT + veth 橋接", "外部拓樸", ["gap"], None, "原始 JSON 無法證明 NAT 路由"),
   ("K-F5", "虛擬 DUT testbench", "非選定配對包", ["gap"], None, "投影資料非 testbench 程式"),
   ("K-F6", "實體 DUT / 流量產生器", "無實體結果保存", ["gap"], None, "不得宣稱實測吞吐 / 流量成功"),
 ]),
 dict(bucket="C", title="報告層交付", status_text="資料缺口（報告層）", status_color=REDST, rows=[
   ("K-F8", "最終書面報告", "報告層另行撰寫", ["gap"], None, "選定配對 HTML 為附錄候選；非最終報告"),
   ("R1-D5", "技術 WP1 評估報告", "可支持未來 WP1，不取代", ["gap"], None, "非完整 WP1 技術評估報告"),
 ]),
]

# ------------------------------------------------------------------ assemble
def s_fig(prs, L, title_main, sub, png, takeaway=None):
    """Full-width figure slide: title bar + aspect-fitted centered image + honesty takeaway line."""
    from PIL import Image
    s = prs.slides.add_slide(L[1])
    tr = [(title_main, {"sz": T_TITLE, "b": True, "c": NAVY})]
    if sub:
        tr.append(("　" + sub, {"sz": T_SUB, "b": True, "c": SOFT}))
    fill_title(s, tr); clear_body_ph(s)
    if not png or not os.path.exists(png):
        return s
    top = TOP_CONTENT
    avail_h = (PAGE_BOTTOM - top) - (0.55 if takeaway else 0.10)
    iw, ih = Image.open(png).size; ar = iw / ih
    w = 12.4; h = w / ar
    if h > avail_h:
        h = avail_h; w = h * ar
    left = (13.333 - w) / 2.0
    s.shapes.add_picture(png, Inches(left), Inches(top), width=Inches(w), height=Inches(h))
    fr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    fr.fill.background(); fr.line.color.rgb = RGBColor(0xD9, 0xDD, 0xE3); fr.line.width = Pt(0.75); fr.shadow.inherit = False
    if takeaway:
        tt = textbox(s, Inches(0.6), Inches(top + h + 0.10), Inches(12.1), Inches(0.45))
        para(tt, [(takeaway, {"sz": 16, "c": NAVY, "b": True})], first=True, sa=0)
    return s

def s_viznative(prs, L, title_main, sub, viz_name):
    """Box/flow figure rendered NATIVELY (real pptx shapes, 標楷體/Times) from its *.viz.json."""
    s = prs.slides.add_slide(L[1])
    tr = [(title_main, {"sz": T_TITLE, "b": True, "c": NAVY})]
    if sub:
        tr.append(("　" + sub, {"sz": T_SUB, "b": True, "c": SOFT}))
    fill_title(s, tr); clear_body_ph(s)
    viz_native.draw_viz(s, os.path.join(FIGDIR, viz_name))
    return s

def main():
    base.build_charts()
    prs = Presentation(TEMPLATE)
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        rid = sid.get(qn("r:id"))
        if rid:
            try: prs.part.drop_rel(rid)
            except Exception: pass
        lst.remove(sid)
    L = prs.slide_layouts

    s_title(prs, L)
    s_agenda(prs, L)
    s_rule(prs, L)
    s_viznative(prs, L, "驗收計分板", "by the numbers", "scoreboard-by-the-numbers.viz.json")
    s_overview(prs, L)
    s_viznative(prs, L, "能力地圖", "白話能力涵蓋 · 依分類與狀態", "capability-map.viz.json")
    s_demo(prs, L)
    s_data_model(prs, L)
    s_ops_tle(prs, L)
    s_ops_cron(prs, L)
    s_ops_gate(prs, L)
    s_viznative(prs, L, "Demo 操作流程", "雙站投影 → 播放 → 證據", "flow-demo-journey.viz.json")
    s_viznative(prs, L, "系統架構與資料流", "公開來源 → 標準模型 → 示意代理", "system-architecture-v2.viz.json")
    s_fig(prs, L, "三軌道高度與週期", "LEO / MEO / GEO 示意", figp("scene-orbit-scale.png"),
          "高度為公開標稱值；半徑為示意壓縮比例（非等比繪製）。")
    s_fig(prs, L, "雙站互視幾何", "互視交集 = 雙站可視時間窗", figp("scene-mutual-visibility.png"),
          "幾何由公開站座標 + SGP4 軌道計算；仰角門檻為模型參數，非實測。")
    s_fig(prs, L, "衛星跨軌道換手場景", "V-MO1 · demo-balanced-v1 · TR 38.821 §7.3.2.2", figp("scene-handover.png"),
          "換手依模型仰角 / 幾何計算，非營運商實測；RSRP 為相對代理。")
    s_fig(prs, L, "鏈路預算沿路徑分解", "FSPL · 雨衰 · 氣體 · 天線（各標標準出處）", figp("scene-link-budget-path.png"),
          "各項依 3GPP / ITU 標準建模；未含 EIRP / 雜訊指數 / 封包實測 / 營運商校準。")
    cov_cols = [Inches(5.95), Inches(1.55), Inches(2.35), Inches(2.5)]
    paginated_table(prs, L, "驗收總覽 · 覆蓋矩陣", "第一類 · 19 列", ["需求", "狀態", "主來源類別", "需求 ID"], cov_rows(base.REQS_A), cov_cols, per_page=10)
    paginated_table(prs, L, "驗收總覽 · 覆蓋矩陣", "第二類 · 7 + 第三類 · 8", ["需求", "狀態", "主來源類別", "需求 ID"], cov_rows(base.REQS_B + base.REQS_C), cov_cols, per_page=10)
    s_viznative(prs, L, "誠實儀表板", "可宣稱 / 代理 / 不可宣稱", "honesty-dashboard.viz.json")
    s_viznative(prs, L, "資料來源 → 模型 → 誠實分級", "provenance & honesty grading", "flow-provenance-honesty.viz.json")
    s_donotclaim(prs, L)
    paginated_table(prs, L, "證據包", "可證 / 不可證", TRACE_HEADERS, TRACE_ROWS,
                    [Inches(4.0), Inches(4.6), Inches(3.9)], per_page=4, bsz=BIG_BODY, hsz=BIG_HEAD)

    s_section(prs, L, "A")
    for g in GROUPS:
        if g["bucket"] == "A":
            s_group(prs, L, g)
    s_section(prs, L, "B")
    s_viznative(prs, L, "來源權威階梯", "現在在 Tier-B，缺一階到 Tier-A", "tier-ladder.viz.json")
    for g in GROUPS:
        if g["bucket"] == "B":
            s_group(prs, L, g)
    s_modelcharts(prs, L)
    s_section(prs, L, "C")
    for g in GROUPS:
        if g["bucket"] == "C":
            s_group(prs, L, g)

    gap_rows = [[(g[0].replace("Bucket B", "第二類"), {"b": True}), g[1], g[2].replace("Bucket B", "第二類")] for g in base.GAP_ROWS]
    paginated_table(prs, L, "資料缺口總表", "已揭露待補（明細）", ["缺口", "阻擋需求", "關閉所需"], gap_rows,
                    [Inches(2.7), Inches(2.9), Inches(6.75)], per_page=4, bsz=BIG_BODY, hsz=BIG_HEAD)
    cat_rows = [[(c[0], {"b": True, "c": NAVY}), c[1], strip_hash(c[2]), (strip_hash(c[3]), {"c": DRED})] for c in base.CATALOG]
    paginated_table(prs, L, "附錄：來源目錄", "[1]–[16]", ["ID", "來源", "可支持", "不可支持"], cat_rows,
                    [Inches(0.7), Inches(3.5), Inches(4.5), Inches(3.65)], per_page=4, bsz=BIG_BODY, hsz=BIG_HEAD)
    s_closing(prs, L)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    # Kill PowerPoint spell-check squiggles: tag EVERY run lang=zh-TW + noProof
    # (PP otherwise proofs generated CJK text against the default en-US dictionary
    # and red-underlines it; copy/paste used to be the only clear-it workaround).
    for _sl in prs.slides:
        for _tag in ("a:rPr", "a:endParaRPr", "a:defRPr"):
            for _rpr in _sl._element.iter(qn(_tag)):
                _rpr.set("lang", "zh-TW")
                _rpr.set("altLang", "zh-TW")  # PP proofs CJK via altLang; en-US here = squiggle
                _rpr.set("noProof", "1")
    prs.save(OUT)
    print(f"slides: {len(prs.slides)}  ->  {OUT}")

if __name__ == "__main__":
    main()
