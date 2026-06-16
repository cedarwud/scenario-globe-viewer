"""Render a research-visual-lab *.viz.json box/flow diagram NATIVELY into a python-pptx
slide (real shapes + text, not an embedded raster). CJK text = 標楷體, Latin = Times New
Roman. Scenes with curved geometry (Earth/orbits) are NOT handled here — keep those as images.

Public: draw_viz(slide, viz_path, top=1.30, bottom=6.9, left=0.45, full_w=12.43)
Maps the viz canvas (minus its 'title' node) into the slide content box, preserving aspect.
"""
import json
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

EA = "標楷體"
LAT = "Times New Roman"


def _setfont(run):
    run.font.name = LAT
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", "zh-TW")
    rPr.set("altLang", "zh-TW")
    rPr.set("noProof", "1")
    for tag in ("a:latin", "a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set("typeface", EA if tag == "a:ea" else LAT)


def _rgb(hexs, default="1a2238"):
    if not hexs or not isinstance(hexs, str) or not hexs.startswith("#"):
        hexs = "#" + default
    h = hexs.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _lum(hexs):
    h = (hexs or "#ffffff").lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    try:
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    except ValueError:
        return 1.0
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255.0


def _anchor(box, side):
    x, y, w, h = box
    return {
        "left": (x, y + h / 2), "right": (x + w, y + h / 2),
        "top": (x + w / 2, y), "bottom": (x + w / 2, y + h),
    }.get(side, (x + w / 2, y + h / 2))


def _arrowhead(shape, size="med"):
    ln = shape.line._get_or_add_ln()
    end = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": size, "len": size})
    ln.append(end)


def draw_viz(slide, viz_path, top=1.32, bottom=6.55, left=0.45, full_w=12.43):
    d = json.load(open(viz_path))
    style = d.get("style", {})
    nstyle = style.get("node", {}); sstyle = style.get("section", {}); estyle = style.get("edge", {})
    astyle = style.get("annotation", {})
    sections = d.get("sections", [])
    nodes = [n for n in d.get("nodes", []) if n.get("id") != "title"]
    annos = d.get("annotations", [])
    edges = d.get("edges", [])
    nbyid = {n["id"]: n for n in d.get("nodes", [])}

    # ---- content bbox over sections + (non-title) nodes + annotations ----
    xs, ys, xe, ye = [], [], [], []
    for s in sections:
        xs.append(s["x"]); ys.append(s["y"]); xe.append(s["x"] + s["w"]); ye.append(s["y"] + s["h"])
    for n in nodes:
        xs.append(n["x"]); ys.append(n["y"]); xe.append(n["x"] + n["w"]); ye.append(n["y"] + n["h"])
    for a in annos:
        xs.append(a["x"]); ys.append(a["y"]); xe.append(a["x"] + a.get("w", 300)); ye.append(a["y"] + 24)
    if not xs:
        return
    minx, miny, maxx, maxy = min(xs), min(ys), max(xe), max(ye)
    bw, bh = maxx - minx, maxy - miny
    avail_h = bottom - top
    sc = min(full_w / bw, avail_h / bh)          # inches per viz-px
    ox = left + (full_w - bw * sc) / 2.0 - minx * sc
    oy = top + (avail_h - bh * sc) / 2.0 - miny * sc

    def X(px): return ox + px * sc
    def Y(py): return oy + py * sc

    def _eff(s):  # effective char-width units: CJK ~1.0em, latin/space ~0.56em
        return sum(1.0 if ord(c) > 0x2E80 else 0.56 for c in s) or 1.0

    def fit_pt(lines, w_in, h_in, target):
        # largest pt <= target such that each label line fits on ONE line (width) and all lines fit (height).
        # target = the viz fontSize used directly as a pt target (body ~15-16, big numbers large -> box-limited).
        if not lines:
            return float(target)
        h_pt = (h_in * 72.0 - 5.0) / max(1, len(lines)) / 1.22
        w_pt = min((w_in * 72.0 - 11.0) / _eff(ln) for ln in lines)
        return max(8.5, min(float(target), h_pt, w_pt))

    def box_in(n):  # node bbox in inches
        return (X(n["x"]), Y(n["y"]), n["w"] * sc, n["h"] * sc)

    # ---- sections (paint first / back) ----
    for s in sections:
        st = {**sstyle, **s.get("style", {})}
        l, t, w, h = X(s["x"]), Y(s["y"]), s["w"] * sc, s["h"] * sc
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        try: shp.adjustments[0] = 0.03
        except Exception: pass
        shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(st.get("fill"), "f7f9fc")
        shp.line.color.rgb = _rgb(st.get("stroke"), "c2cede"); shp.line.width = Pt(1.0)
        shp.shadow.inherit = False
        shp.text_frame.word_wrap = True
        lbl = s.get("label", "")
        if lbl:
            tb = slide.shapes.add_textbox(Inches(l + 0.12), Inches(t + 0.04), Inches(w - 0.24), Inches(0.34))
            p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = lbl
            slbl_pt = min(13.0, max(10.0, ((w - 0.24) * 72.0 - 6.0) / (_eff(lbl) or 1.0)))
            r.font.size = Pt(slbl_pt); r.font.bold = True
            r.font.color.rgb = _rgb(st.get("textFill"), "44506a"); _setfont(r)
            tb.text_frame.word_wrap = True

    # ---- edges (over sections, under nodes) ----
    for e in edges:
        a, b = nbyid.get(e.get("from")), nbyid.get(e.get("to"))
        if not a or not b:
            continue
        est = {**estyle, **e.get("style", {})}
        p1 = _anchor((a["x"], a["y"], a["w"], a["h"]), e.get("fromAnchor"))
        p2 = _anchor((b["x"], b["y"], b["w"], b["h"]), e.get("toAnchor"))
        wps = [(w["x"], w["y"]) for w in e.get("waypoints", [])]
        if not wps and e.get("route") == "orthogonal":
            if e.get("fromAnchor") in ("left", "right"):
                midx = (p1[0] + p2[0]) / 2.0; wps = [(midx, p1[1]), (midx, p2[1])]
            else:
                midy = (p1[1] + p2[1]) / 2.0; wps = [(p1[0], midy), (p2[0], midy)]
        pts = [p1] + wps + [p2]
        emu = [(Emu(int(Inches(X(px)))), Emu(int(Inches(Y(py))))) for px, py in pts]
        fb = slide.shapes.build_freeform(emu[0][0], emu[0][1], scale=1.0)
        fb.add_line_segments(emu[1:], close=False)
        shp = fb.convert_to_shape()
        shp.fill.background()
        shp.line.color.rgb = _rgb(est.get("stroke"), "2a9d8f")
        shp.line.width = Pt(max(1.0, est.get("strokeWidth", 2.0) * 0.85))
        if est.get("dash"):
            shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        if est.get("arrowSize", 11) > 0:
            _arrowhead(shp)
        shp.shadow.inherit = False

    # ---- nodes (front) ----
    for n in nodes:
        ns = {**nstyle, **n.get("style", {})}
        l, t, w, h = box_in(n)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        try: shp.adjustments[0] = 0.08
        except Exception: pass
        fill = ns.get("fill", "#ffffff"); shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(fill, "ffffff")
        shp.line.color.rgb = _rgb(ns.get("stroke"), "1a2238")
        shp.line.width = Pt(max(0.75, ns.get("strokeWidth", 1.5) * 0.7))
        shp.shadow.inherit = False
        tf = shp.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for m in ("margin_left", "margin_right"): setattr(tf, m, Inches(0.04))
        for m in ("margin_top", "margin_bottom"): setattr(tf, m, Inches(0.02))
        dark = _lum(fill) < 0.45
        tcol = _rgb(ns.get("textFill"), "ffffff" if dark else "1a2238")
        lines = n.get("label", [])
        if isinstance(lines, str): lines = [lines]
        node_pt = fit_pt(lines, w, h, ns.get("fontSize", 16))
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = ln
            r.font.size = Pt(node_pt)
            r.font.bold = bool(dark)
            r.font.color.rgb = tcol; _setfont(r)

    # ---- annotations (footers, front) ----
    for a in annos:
        ast = {**astyle, **a.get("style", {})}
        l, t = X(a["x"]), Y(a["y"]); w = a.get("w", 1000) * sc
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(max(2.0, w)), Inches(0.66))
        tf = tb.text_frame; tf.word_wrap = True
        # keep footers compact so they fit one line where possible (avoid spill into the template footer)
        afit = min(12.0, max(9.5, (w * 72.0 - 6.0) / (_eff(max(a.get("label", [""]), key=len)) or 1.0)))
        lines = a.get("label", [])
        if isinstance(lines, str): lines = [lines]
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = ln
            r.font.size = Pt(afit)
            r.font.color.rgb = _rgb(ast.get("textFill"), "6b7280"); _setfont(r)
