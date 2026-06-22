#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SAMPLE friendly slides (style test v2) — capability cards + 三層燈號 honesty scale
(lamp at the page footer, right of the wMNLab logo) + speaker notes (被問時答法).
Professional noun-phrase labels (not colloquial questions); fuller explanations.
Content top aligned with the main deck. Code-verified content.

Output: deliverable/sample-friendly-slides.pptx
Run: python3 scripts/build_sample_friendly_slides.py
"""
import os, sys
HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import build_requirement_presentation_v2 as v2

ROOT = os.path.dirname(HERE)
TEMPLATE = os.path.join(ROOT, "template.pptx")
OUT = os.path.join(ROOT, "deliverable", "sample-friendly-slides.pptx")

NAVY, INK, SOFT, DRED = v2.NAVY, v2.INK, v2.SOFT, v2.DRED
GREEN, AMBER, REDST = v2.GREEN, v2.AMBER, v2.REDST
setfont, para, fit, fill_title, clear_body_ph, textbox = (
    v2.setfont, v2.para, v2.fit, v2.fill_title, v2.clear_body_ph, v2.textbox)

TOP = 1.05   # content top, raised ~1 font-size higher (verify clears the title rule)
TIER = {"g": (GREEN, "已實作展示"), "y": (AMBER, "標準公式模型"), "r": (REDST, "需營運商實測")}

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def s_legend(prs, L):
    s = prs.slides.add_slide(L[1])
    fill_title(s, [("資料誠實分級 · 三層燈號", {"sz": 28, "b": True, "c": NAVY}),
                   ("　燈號標於每頁頁尾（wMNLab 標誌右側）", {"sz": 16, "b": True, "c": SOFT})])
    clear_body_ph(s)
    tf = textbox(s, Inches(0.7), Inches(TOP), Inches(12.0), Inches(5.75), shrink=True)
    rows = [
        ("g", "功能已實作,於 demo 畫面可實際操作與檢視。", "視覺化 · 互動 UI · 報告匯出"),
        ("y", "依 ITU-R / 3GPP 國際標準公式即時計算所得之模型值,非實測。", "傳播延遲 · 降雨衰減 · 換手決策 · 天線增益"),
        ("r", "需營運商實測資料方能成立,目前尚無,已誠實標註為「待補」。", "封包測試 · 原生 RF 換手 · 站台硬體實測"),
    ]
    for i, (k, desc, eg) in enumerate(rows):
        col, name = TIER[k]
        para(tf, [("● ", {"sz": 28, "c": col}), (name + "　", {"sz": 22, "b": True, "c": col}), (desc, {"sz": 20, "c": INK})],
             first=(i == 0), sa=4, line=1.22)
        para(tf, [("　　例：" + eg, {"sz": 14, "c": SOFT})], sa=16, line=1.1)
    para(tf, [("使用方式：每頁標示其對應燈號;受詢問時,先依燈號層級,再援引該頁「依據 / 邊界」一句說明。",
               {"sz": 16, "b": True, "c": NAVY})], sa=0, line=1.22)
    notes(s, "口白：本簡報所有畫面數據均標示三層燈號。綠燈為已實作、可操作展示;黃燈為依國際標準公式即時計算之模型值,非實測;"
             "紅燈為需營運商實測、目前尚無、已標待補。受詢問『是否屬實測』時,先依燈號層級回應,再援引該頁依據與邊界。"
             "誠實原則:模型 / 公式 / 公開推導之值,不陳述為實測或營運商驗證。")
    return s

def card(prs, L, title, tier, points, basis, boundary, reqs, note_text):
    s = prs.slides.add_slide(L[1])
    fill_title(s, [(title, {"sz": 26, "b": True, "c": NAVY})])
    clear_body_ph(s)
    col, name = TIER[tier]
    # body (raised to standard content top)
    tf = textbox(s, Inches(0.7), Inches(TOP), Inches(11.9), Inches(5.75), shrink=True)
    first = True
    for label, txt in points:
        para(tf, [(label + "　", {"sz": 20, "b": True, "c": NAVY}), (txt, {"sz": 20, "c": INK})], first=first, sa=10, line=1.26)
        first = False
    para(tf, [("依據　", {"sz": 16, "b": True, "c": SOFT}), (basis, {"sz": 16, "c": SOFT})], sb=5, sa=4, line=1.18)
    para(tf, [("邊界　", {"sz": 16, "b": True, "c": DRED}), (boundary, {"sz": 16, "c": DRED})], sa=0, line=1.18)
    # footer: tier lamp (right of wMNLab logo) + requirement tags (right)
    lamp = textbox(s, Inches(2.6), Inches(7.0), Inches(4.3), Inches(0.34), anchor=MSO_ANCHOR.MIDDLE)
    para(lamp, [("● ", {"sz": 16, "c": col}), (name, {"sz": 13, "b": True, "c": col})], first=True, sa=0)
    ft = textbox(s, Inches(7.0), Inches(7.0), Inches(5.4), Inches(0.34), anchor=MSO_ANCHOR.MIDDLE)
    para(ft, [("對應需求　" + reqs, {"sz": 12, "c": SOFT})], align=PP_ALIGN.RIGHT, first=True, sa=0)
    notes(s, note_text)
    return s

def main():
    prs = Presentation(TEMPLATE)
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        rid = sid.get(v2.qn("r:id"))
        if rid:
            try: prs.part.drop_rel(rid)
            except Exception: pass
        lst.remove(sid)
    L = prs.slide_layouts

    s_legend(prs, L)

    card(prs, L, "衛星資料：來源、篩選與運算", "g",
         [("資料來源", "採 CelesTrak 公開 GP/TLE 軌道資料,專案內已打包為本地檔,共 1258 顆（OneWeb 651 顆 LEO、Galileo 33 顆 MEO、CelesTrak GEO 群組 574 顆）,最新快照 2026-06-14;瀏覽器執行時讀取本地檔,不連外網下載。"),
          ("篩選條件", "以 SGP4 推算各衛星對兩站之可見時間窗,取「兩站同時可見」之交集;站點仰角須高於有效門檻（基礎 10° 加地形遮罩,CHT 為 31°、SANSA 為 11°）,各軌道類別再依互動運算預算截頂（LEO 200 / MEO 100 / GEO 60）。"),
          ("運算方式", "採瀏覽器端即時運算 — 以 satellite.js（SGP4 演算法）針對本配對與選定時間窗當場推算衛星位置與可見性,並非載入預先計算之結果檔。"),
          ("可見性判定", "畫面僅呈現「兩站幾何上互相可見」之候選衛星,為純幾何推算結果,非營運商宣稱或預設清單。")],
         "CelesTrak GP/TLE · satellite.js v7（SGP4）· 仰角門檻 10° 加地形遮罩",
         "快照非即時下載最新 TLE;幾何可見不等於已建立連線或 RF 量測。",
         "R1-T1 · K-A1 · R1-F1 · K-A4",
         "受詢問時：\n"
         "Q 衛星資料來源? → CelesTrak 公開 TLE,已打包於專案內,共 1258 顆（最新快照 2026-06-13）。\n"
         "Q 是否即時運算? → 是,瀏覽器以 satellite.js 之 SGP4 即時推算;惟採用打包之快照 TLE,非即時下載最新目錄。\n"
         "Q 是否確為兩站可見之衛星? → 是,僅呈現兩站幾何互可見之交集;惟此為幾何可見,不代表已建立連線。\n"
         "Q 為何非全部衛星? → 來源池 OneWeb 651 LEO、Galileo 33 MEO、GEO 群組 574;再依互動運算預算每類截頂（LEO 200 / MEO 100 / GEO 60），並只取兩站互可見者。")

    card(prs, L, "降雨衰減：模型來源、參數意義與頻段差異", "y",
         [("模型來源", "採 ITU-R P.618-14 §2.2.1.1 完整 slant-path 雨衰預測法,搭配 ITU-R P.838-3 之 k / α 比衰減係數表。"),
          ("降雨參數意義", "降雨輸入代表「瞬時降雨率（mm/h）」之假設情境（what-if）;數值提高將增加雨衰（dB）,連帶降低吞吐估計並加大抖動。demo 預設 5 mm/h,係由公開 CWA 觀測導出之保守值。"),
          ("MEO 不受影響之原因", "本模型之雨衰計算僅於 10–30 GHz 載波區間啟用;MEO 採 L-band 1.5 GHz,低於下限,故雨衰設為 0。此與物理特性一致 — L-band 雨衰趨近於 0,LEO（Ku 12 GHz）與 GEO（Ka 20 GHz）則影響顯著。")],
         "ITU-R P.618-14 · ITU-R P.838-3 · 載波 LEO 12 / MEO 1.5 / GEO 20 GHz",
         "屬 what-if 假設雨率,非當地實測天氣,亦非 R0.01 長期統計校正。",
         "K-E6 · K-A3-b",
         "受詢問時：\n"
         "Q 雨衰公式來源? → ITU-R P.618-14 國際標準路徑法,係數採 P.838-3。\n"
         "Q 降雨參數之意義? → 假設之瞬時降雨率（mm/h）,為 what-if 情境,非實測天氣;預設 5 mm/h 由公開 CWA 觀測導出。\n"
         "Q MEO 為何不受影響? → MEO 載波為 L-band 1.5 GHz,本模型雨衰僅於 10–30 GHz 計算,故設為 0;且物理上 L-band 雨衰本即趨近 0,LEO / GEO 之 Ku / Ka 高頻方才顯著。")

    card(prs, L, "換手：性質、判斷條件與訊號依據", "y",
         [("換手性質", "依 3GPP TR 38.821 §7.3.2.2（條件式換手 CHO）啟發之政策,於候選衛星間擇優;服務對象變更即記為一次換手事件。demo route 採 demo-balanced-v1 跨軌輪替策略（cross-orbit-live 為通用 default）;下列判斷條件為通用跨軌邏輯之概念說明,而 CHT×SANSA demo 之 LEO 互視窗為 0,實際呈現 GEO⇄MEO。"),
          ("跨軌判斷條件", "須同時滿足三項：（一）該時刻兩軌道皆存在「兩站互可見」之衛星;（二）現役 LEO 仰角低於門檻,且無其他 LEO 達標;（三）存在達標之 MEO / GEO 候選,且其訊號（代理 RSRP）高出遲滯（hysteresis）門檻。"),
          ("達標定義", "仰角不低於門檻、剩餘可見時間充足、且傳播延遲在預算之內。"),
          ("訊號依據", "候選比較採相對 RSRP 代理值（天線增益為假設 Tier-B 參數、EIRP 未知）,僅供排序使用,不代表絕對接收功率。")],
         "3GPP TR 38.821 §7.3.2.2 · V-MO1 · 仰角 + 可見時間 + 遲滯門檻",
         "屬模型政策之服務對象改選,非原生 RF 換手,亦非真實服務路由遷移。",
         "R1-T5 · R1-F4 · K-E4 · V-MO1",
         "受詢問時：\n"
         "Q 是否為原生 RF 換手? → 否。係模型依政策於候選衛星間擇優改選;真實 RF 換手須營運商或終端之事件日誌。\n"
         "Q 如何判斷該時刻可跨軌換手? → 程式逐時檢查:兩軌皆有互可見衛星、現役 LEO 低於仰角門檻、且存在訊號高出遲滯門檻之 MEO / GEO 候選,三者皆成立方觸發。\n"
         "Q RSRP 是否為實測? → 否,為相對代理值（天線增益採假設 Tier-B、EIRP 未知）,僅供候選排序比較。")

    card(prs, L, "站點資料：來源、能力依據與真實性", "y",
         [("資料來源", "採公開來源人工策展之站台登錄表,涵蓋營運商官方頁面、FCC / ITU 申報文件、Wikipedia 及 World Teleport Association 等。"),
          ("能力判斷依據", "軌道支援能力取自登錄表之 supportedOrbits 欄位（公開宣稱）;配對可信度則採保守認定 — 目前登錄表無任何配對級背書（pairAttestation）,故所有配對均歸為「幾何推導」,同營運商家族亦不予採認。"),
          ("座標與高程真實性", "CHT Yangmingshan 為代表性區域座標,非精確設施點;高程 489 m 與地形遮罩係公開 Copernicus DEM GLO-30 衍生值,非營運商實測;SANSA 高程 1553 m 為營運商宣稱值。")],
         "multi-orbit-public-registry（公開策展）· Copernicus DEM GLO-30",
         "公開宣稱能力不等於該時段必有三軌;非營運商私有驗證,亦非精確設施座標。",
         "R1-T2 · K-D1",
         "受詢問時：\n"
         "Q 站點資料來源? → 公開來源策展:營運商官方頁面、FCC / ITU 申報、Wikipedia、WTA。\n"
         "Q 以何依據判斷站點具三軌能力? → 取自登錄表之公開宣稱能力欄位;惟對「配對」之可信度採保守認定,目前所有配對皆為幾何推導,無任何配對級公開背書。\n"
         "Q 座標與高程是否屬實測? → CHT 為代表性區域座標,非精確設施點;高程與地形遮罩為公開 Copernicus DEM 衍生,非營運商實測,均已誠實標註。")

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"slides: {len(prs.slides)}  ->  {OUT}")

if __name__ == "__main__":
    main()
